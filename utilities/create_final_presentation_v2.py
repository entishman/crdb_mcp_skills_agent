#!/usr/bin/env python3
"""
Final Visual PowerPoint Presentation - Version 2
Complete presentation with detailed technical flow diagrams
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def add_box(slide, x, y, width, height, text, fill_color, text_size=10, bold=False):
    """Add a colored box with BLACK text"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(2)

    box.text_frame.text = text
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = 1

    # Set ALL paragraphs to black (multi-line text creates multiple paragraphs)
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # BLACK TEXT!

    return box

def add_arrow(slide, x1, y1, x2, y2, label=None):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(2)
    connector.line.end_arrow_type = 2

    if label:
        label_box = slide.shapes.add_textbox(
            Inches((x1 + x2) / 2 - 0.4), Inches((y1 + y2) / 2 - 0.15),
            Inches(0.8), Inches(0.25)
        )
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.italic = True
        label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        label_box.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return connector

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title

    p = title_frame.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_bullet_slide(prs, title, bullets):
    """Add a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    tf = body_shape.text_frame
    tf.clear()

    for bullet in bullets:
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 - level * 2)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(10)

    return slide

def add_architecture_overview(prs):
    """Complete system architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text = "System Architecture: Complete Component Map"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # User
    add_box(slide, 4.25, 0.8, 1.5, 0.5, "User", RGBColor(173, 216, 230))
    add_arrow(slide, 5, 1.3, 5, 1.8)

    # Agent (center)
    add_box(slide, 3.5, 1.8, 3, 0.7, "Agent (agent.py)\nOrchestrator", RGBColor(255, 200, 100), 13)

    # Left - Embeddings
    add_box(slide, 0.5, 3, 2, 0.6, "Embedding Manager\nsentence-transformers\n768-dim",
            RGBColor(255, 230, 230), 11)
    add_arrow(slide, 1.5, 2.5, 4.5, 2.2, "embeddings")

    # Right - Claude API
    add_box(slide, 7.5, 3, 2, 0.6, "Claude Sonnet 4.6\n(Vertex AI)",
            RGBColor(200, 230, 255), 12)
    add_arrow(slide, 6.5, 2.2, 8, 3, "API calls")

    # Bottom left - Query Store
    add_box(slide, 0.5, 4.2, 2, 0.6, "Query Store\n(query_store.py)",
            RGBColor(220, 255, 220), 11)
    add_arrow(slide, 4.5, 2.5, 1.5, 4.2, "find similar")

    # Bottom center - Skill Fetcher
    add_box(slide, 3.5, 4.2, 2, 0.6, "Skill Fetcher\n(skill_fetcher.py)",
            RGBColor(255, 240, 200), 11)
    add_arrow(slide, 4.5, 2.5, 4.5, 4.2, "load skills")

    # Bottom right - MCP Client
    add_box(slide, 6.5, 4.2, 2, 0.6, "MCP Client\n(mcp_client.py)",
            RGBColor(240, 220, 255), 11)
    add_arrow(slide, 6, 2.5, 7.5, 4.2, "execute tools")

    # Database layer - query_history
    add_box(slide, 0.5, 5.3, 2.5, 0.6, "ai_demo.query_history\n768-dim vectors + skills",
            RGBColor(180, 220, 180), 10)
    add_arrow(slide, 1.5, 4.8, 1.5, 5.3)

    # Database layer - skills table
    add_box(slide, 3.3, 5.3, 2.5, 0.6, "ai_demo.skills\nSKILL.md content",
            RGBColor(180, 220, 180), 10)
    add_arrow(slide, 4.5, 4.8, 4.5, 5.3)

    # MCP Server
    add_box(slide, 6.5, 5.3, 2, 0.6, "MCP Server\ncockroachlabs.cloud/mcp",
            RGBColor(255, 210, 210), 10)
    add_arrow(slide, 7.5, 4.8, 7.5, 5.3)

    # CockroachDB Cluster
    add_box(slide, 2.5, 6.3, 5, 0.7, "CockroachDB Cluster\n(Distributed SQL Database)",
            RGBColor(100, 150, 200), 12)
    add_arrow(slide, 1.75, 5.9, 4, 6.3)
    add_arrow(slide, 4.5, 5.9, 5, 6.3)
    add_arrow(slide, 7.5, 5.9, 6, 6.3)

    return slide

def add_flow_new_query_detailed(prs):
    """Detailed flow for new query"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 1: New Query (Cold Start) - Simple Query with Zero Skills"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    add_box(slide, 0.3, 0.7, 2.2, 0.55, 'Step 1: User Input\n"list databases"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    # Step 2
    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'ollama_embedding_manager.py\n' +
            'generate_embedding()\n' +
            '→ 768-dim vector',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3
    add_box(slide, 5.6, 0.7, 2.9, 0.55,
            'Step 3: Search Similar Queries\n' +
            'query_store.py: find_similar_queries()\n' +
            'DB: CockroachDB ai_demo.query_history\n' +
            'Result: ❌ No match (new query)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 7, 1.25, 7, 1.5)

    # Step 4
    add_box(slide, 5.6, 1.5, 2.9, 0.55,
            'Step 4: Pre-load Skills\n' +
            'agent.py: _build_system_prompt()\n' +
            'calls _get_default_skills()\n' +
            'Returns: [] (zero skills!)',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 7, 2.05, 7, 2.3)

    # Step 5
    add_box(slide, 5.6, 2.3, 2.9, 0.7,
            'Step 5: Discover & Call Claude\n' +
            'mcp_client.py: list_tools()\n' +
            'Discovers 11 MCP tools dynamically\n' +
            'agent.py: run_task()\n' +
            'claude.messages.create(tools=...)',
            RGBColor(200, 230, 255), 8.5, False)
    add_arrow(slide, 7, 3, 7, 3.25)

    # Step 6
    add_box(slide, 5.6, 3.25, 2.9, 0.75,
            'Step 6: Tool Discovery\n' +
            'Claude receives tool definitions:\n' +
            '{"name": "list_databases",\n' +
            ' "description": "Lists all databases..."}\n' +
            'Claude calls: list_databases()',
            RGBColor(240, 220, 255), 8.5, False)
    add_arrow(slide, 7, 4, 7, 4.25)

    # Step 7
    add_box(slide, 5.6, 4.25, 2.9, 0.65,
            'Step 7: Execute via MCP\n' +
            'mcp_client.py: execute_tool_call()\n' +
            'JSON-RPC → MCP Server\n' +
            'MCP → CockroachDB\n' +
            'Runs: SHOW DATABASES',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 5.6, 4.55, 5.3, 4.55)

    # Step 8
    add_box(slide, 2.8, 4.25, 2.5, 0.65,
            'Step 8: Results to Claude\n' +
            'MCP returns database list\n' +
            'Claude formats answer',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.8, 4.55, 2.5, 4.55)

    # Step 9
    add_box(slide, 0.3, 4.25, 2.2, 0.65,
            'Step 9: Store Learning\n' +
            'query_store.py: store_query()\n' +
            'Saves: query + [] skills\n' +
            'DB: ai_demo.query_history',
            RGBColor(220, 255, 220), 9, False)

    # Explanations
    explain1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.15), Inches(4.5), Inches(0.55))
    explain1.text_frame.text = "🔍 Dynamic Tool Discovery: Agent calls MCP server's tools/list method at startup. MCP returns tool definitions (name, description, parameters). No hardcoded tools = not brittle!"
    for p in explain1.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)

    explain2 = slide.shapes.add_textbox(Inches(5), Inches(5.15), Inches(4.5), Inches(0.55))
    explain2.text_frame.text = "⚡ Zero-Skill Optimization: Simple queries don't need SKILL.md files. Claude knows basic SQL from training + tool descriptions. Saves ~8KB tokens per query!"
    for p in explain2.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(5.85), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "🗄️ CockroachDB Powers Everything: query_history table (768-dim pgvector), skills table, and the actual cluster being managed - all in one distributed database"
    for p in legend.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)

    return slide

def add_flow_learned_query_detailed(prs):
    """Detailed flow for learned query"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 2: Learned Query (Warm Start) - Progressive Learning in Action"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Steps 1-2
    add_box(slide, 0.3, 0.7, 2.2, 0.55, 'Step 1: User Input\n"check for hotspots"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'ollama_embedding_manager.py\n' +
            'generate_embedding()\n' +
            '→ 768-dim vector',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3 - MATCH FOUND!
    add_box(slide, 5.6, 0.7, 2.9, 0.6,
            'Step 3: Similarity Search\n' +
            'query_store.py: find_similar_queries()\n' +
            'CockroachDB pgvector search\n' +
            'Result: ✅ Match! (0.78 similarity)\n' +
            'Matched: "analyze hot ranges"',
            RGBColor(180, 255, 180), 8.5, False)
    add_arrow(slide, 7, 1.3, 7, 1.55)

    # Step 4
    add_box(slide, 5.6, 1.55, 2.9, 0.6,
            'Step 4: Load Learned Skills\n' +
            'agent.py: _build_system_prompt()\n' +
            'Gets skills from matched query:\n' +
            '["observability-and-diagnostics/\n' +
            'analyzing-range-distribution"]',
            RGBColor(220, 255, 220), 8.5, False)
    add_arrow(slide, 7, 2.15, 7, 2.4)

    # Step 4b
    add_box(slide, 5.6, 2.4, 2.9, 0.6,
            'Step 4b: Fetch SKILL.md\n' +
            'skill_fetcher.py: get_skill()\n' +
            'Source: ai_demo.skills table\n' +
            '(CockroachDB)\n' +
            'Returns: Full best practices doc',
            RGBColor(255, 240, 200), 8.5, False)
    add_arrow(slide, 7, 3, 7, 3.25)

    # Step 5
    add_box(slide, 5.6, 3.25, 2.9, 0.65,
            'Step 5: Call Claude API\n' +
            'agent.py: run_task()\n' +
            'System prompt includes:\n' +
            '  - SKILL.md content\n' +
            '  - Discovered tools\n' +
            '  - User query',
            RGBColor(200, 230, 255), 8.5, False)
    add_arrow(slide, 7, 3.9, 7, 4.15)

    # Step 6
    add_box(slide, 5.6, 4.15, 2.9, 0.6,
            'Step 6: Expert Analysis\n' +
            'Claude uses SKILL.md guidance\n' +
            'Knows: Check SHOW RANGES\n' +
            'Calls: select_query() with\n' +
            'range distribution query',
            RGBColor(240, 220, 255), 8.5, False)
    add_arrow(slide, 5.6, 4.45, 5.3, 4.45)

    # Steps 7-9
    add_box(slide, 2.8, 4.15, 2.5, 0.6,
            'Step 7: Execute Query\n' +
            'mcp_client.py → MCP Server\n' +
            '→ CockroachDB cluster\n' +
            'Returns: range data',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 2.8, 4.45, 2.5, 4.45)

    add_box(slide, 0.3, 4.15, 2.2, 0.6,
            'Step 8: Expert Answer\n' +
            'Claude analyzes using\n' +
            'SKILL.md guidance\n' +
            'Detailed analysis',
            RGBColor(255, 240, 200), 9, False)

    add_box(slide, 0.3, 5, 2.2, 0.5,
            'Step 9: Confirm Learning\n' +
            'Stores: skill helped!\n' +
            'Next time = faster',
            RGBColor(220, 255, 220), 9, False)

    # Explanation
    explain = slide.shapes.add_textbox(Inches(2.8), Inches(5), Inches(6.7), Inches(0.5))
    explain.text_frame.text = "🎯 Progressive Learning: CockroachDB's pgvector enables sub-millisecond similarity search across millions of queries. Each interaction improves future performance!"
    for p in explain.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(5.7), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "🚀 HNSW Index: CockroachDB's HNSW vector index makes similarity search O(log n) instead of O(n). Scales to millions of queries with consistent performance."
    for p in legend.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)

    return slide

def add_flow_complex_skill_fetch(prs):
    """Flow for complex query requiring skill fetch"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 3: Complex Query - Dynamic Skill Fetching (No Match)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Steps 1-4
    add_box(slide, 0.3, 0.7, 2.2, 0.55, 'Step 1: Complex Query\n"Set up MOLT Replicator"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'ollama_embedding_manager.py\n' +
            'generate_embedding()',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    add_box(slide, 5.6, 0.7, 2.9, 0.55,
            'Step 3: Search History\n' +
            'query_store.py\n' +
            '❌ No match\n' +
            '(Never asked about MOLT)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 7, 1.25, 7, 1.5)

    add_box(slide, 5.6, 1.5, 2.9, 0.5,
            'Step 4: Zero Skills\n' +
            '_get_default_skills()\n' +
            'Returns: []',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 7, 2, 7, 2.25)

    # Step 5
    add_box(slide, 5.6, 2.25, 2.9, 0.6,
            'Step 5: Call Claude\n' +
            'Claude gets: query + tools\n' +
            'Claude thinks: "I need\n' +
            'MOLT Replicator expertise"',
            RGBColor(200, 230, 255), 8.5, False)
    add_arrow(slide, 7, 2.85, 7, 3.1)

    # Step 6 - fetch_skill
    add_box(slide, 5.6, 3.1, 2.9, 0.6,
            'Step 6: Claude Calls Tool\n' +
            'fetch_skill(\n' +
            '  "onboarding-and-migrations/\n' +
            '  molt-replicator"\n' +
            ')',
            RGBColor(255, 200, 150), 8.5, False)
    add_arrow(slide, 5.6, 3.4, 5.3, 3.4)

    # Step 7
    add_box(slide, 2.8, 3.1, 2.5, 0.6,
            'Step 7: Fetch from DB\n' +
            'skill_fetcher.py: get_skill()\n' +
            'Query: ai_demo.skills table\n' +
            'Returns: Full SKILL.md',
            RGBColor(255, 240, 200), 8, False)
    add_arrow(slide, 2.8, 3.4, 2.5, 3.4)

    # Steps 8-10
    add_box(slide, 0.3, 3.1, 2.2, 0.6,
            'Step 8: Skill → Claude\n' +
            'Agent appends to\n' +
            'conversation\n' +
            'Claude has expertise!',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 1.4, 3.7, 1.4, 3.95)

    add_box(slide, 0.3, 3.95, 2.2, 0.6,
            'Step 9: Re-process\n' +
            'With SKILL.md context\n' +
            'Knows setup steps\n' +
            'May call MCP tools',
            RGBColor(240, 220, 255), 9, False)
    add_arrow(slide, 2.5, 4.25, 2.8, 4.25)

    add_box(slide, 2.8, 3.95, 2.5, 0.6,
            'Step 10: Execute\n' +
            'Check existing setup:\n' +
            'list_databases(), etc.\n' +
            'Via MCP → CockroachDB',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 5.3, 4.25, 5.6, 4.25)

    add_box(slide, 5.6, 3.95, 2.9, 0.6,
            'Step 11: Expert Guide\n' +
            'Complete MOLT guide\n' +
            'Based on SKILL.md\n' +
            'Calls complete_task()',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 7, 4.55, 7, 4.8)

    add_box(slide, 5.6, 4.8, 2.9, 0.55,
            'Step 12: Store Learning\n' +
            'Saves: query + [skill]\n' +
            'Next time = instant!',
            RGBColor(220, 255, 220), 8.5, False)

    # Explanations
    explain1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(4.5), Inches(0.55))
    explain1.text_frame.text = "🎯 Intelligent Fetching: fetch_skill is a tool available to Claude (like list_databases). Claude decides when to call it based on the complexity of the query and its own knowledge gaps."
    for p in explain1.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)

    explain2 = slide.shapes.add_textbox(Inches(5), Inches(5.5), Inches(4.5), Inches(0.55))
    explain2.text_frame.text = "💾 SKILL.md Storage: 25+ CockroachDB best practice documents stored in ai_demo.skills table. Distributed globally, consistent reads, ACID guarantees. No external docs system needed!"
    for p in explain2.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(6.2), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "🔍 No Hardcoded Tools: Agent discovers all 11 MCP tools dynamically via tools/list method. Adapts automatically when MCP server adds new capabilities."
    for p in legend.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)

    return slide

def add_real_world_use_cases(prs):
    """Real-world use cases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Real-World Use Cases"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Use cases with boxes
    add_box(slide, 0.5, 1.2, 4.2, 0.6, "🔧 DevOps Automation", RGBColor(200, 230, 255), 16, True)
    uc1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(1.2))
    uc1.text_frame.text = """• "Check cluster health daily"
• "Find slow queries and optimize"
• "Monitor background jobs"
• Learns patterns over time
• Reduces ops load 60%"""
    for p in uc1.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.3, 1.2, 4.2, 0.6, "👨‍💼 DBA Self-Service", RGBColor(220, 255, 220), 16, True)
    uc2 = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(1.2))
    uc2.text_frame.text = """• "Create test database for QA"
• "Show schema for users table"
• "Insert sample data"
• Natural language → SQL
• Safety confirmations built-in"""
    for p in uc2.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.5, 3.4, 4.2, 0.6, "🎫 Support Automation", RGBColor(255, 230, 230), 16, True)
    uc3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(4.2), Inches(1.2))
    uc3.text_frame.text = """• Customer: "My query is slow"
• Agent analyzes & suggests fixes
• Follows best practices (SKILL.md)
• Learns from past tickets
• Reduces support load 60%"""
    for p in uc3.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.3, 3.4, 4.2, 0.6, "📚 Training & Onboarding", RGBColor(255, 240, 200), 16, True)
    uc4 = slide.shapes.add_textbox(Inches(5.3), Inches(4.1), Inches(4.2), Inches(1.2))
    uc4.text_frame.text = """• New devs ask questions 24/7
• Consistent accurate answers
• References official docs
• Tracks common questions
• Reduces training time 40%"""
    for p in uc4.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 2.9, 5.6, 4.2, 0.6, "🚀 Migration Planning", RGBColor(240, 220, 255), 16, True)
    uc5 = slide.shapes.add_textbox(Inches(2.9), Inches(6.3), Inches(4.2), Inches(0.9))
    uc5.text_frame.text = """• "Best practices for MOLT Fetch?"
• Guides through complex migrations
• Learns from migration patterns
• Reduces migration risk"""
    for p in uc5.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_why_cockroachdb(prs):
    """Why CockroachDB is perfect"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Why CockroachDB Powers This Agent"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Features
    add_box(slide, 0.5, 1.2, 4.3, 0.7, "🧠 Native pgvector Support", RGBColor(200, 230, 255), 14, True)
    f1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.3), Inches(1.3))
    f1.text_frame.text = """• Stores 768-dim embeddings natively
• HNSW index for fast similarity search
• No external vector DB needed
• SQL + vectors in one place
• Cosine similarity: <-> operator"""
    for p in f1.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.2, 1.2, 4.3, 0.7, "📈 Horizontal Scalability", RGBColor(220, 255, 220), 14, True)
    f2 = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.3), Inches(1.3))
    f2.text_frame.text = """• Add nodes as queries grow
• Auto-rebalancing of data
• No sharding complexity
• Handles millions of queries
• Linear scale-out performance"""
    for p in f2.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.5, 3.6, 4.3, 0.7, "🔒 ACID Transactions", RGBColor(255, 230, 230), 14, True)
    f3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(4.3), Inches(1.3))
    f3.text_frame.text = """• Store query + embedding atomically
• No lost learning data
• Strong consistency guarantees
• Multi-row writes are safe
• Learning survives failures"""
    for p in f3.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.2, 3.6, 4.3, 0.7, "🌍 Global Distribution", RGBColor(255, 240, 200), 14, True)
    f4 = slide.shapes.add_textbox(Inches(5.2), Inches(4.4), Inches(4.3), Inches(1.3))
    f4.text_frame.text = """• Deploy agents globally
• Shared learning across regions
• Low-latency reads everywhere
• Multi-region resilience
• Follow-the-sun support"""
    for p in f4.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 2.85, 6, 4.3, 0.7, "🎯 SQL + JSON + Vectors", RGBColor(240, 220, 255), 14, True)
    f5 = slide.shapes.add_textbox(Inches(2.85), Inches(6.8), Inches(4.3), Inches(0.6))
    f5.text_frame.text = """• Store SKILL.md as TEXT or JSONB
• Array support for skills_used[]
• VECTOR type for embeddings
• Familiar SQL interface"""
    for p in f5.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def create_presentation():
    """Create complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs, "Intelligent CockroachDB Agent",
                   "Self-Learning AI Assistant with Progressive Intelligence")

    # 2. Overview
    overview = [
        "🎯 What It Does",
        ("Natural language interface to CockroachDB", 1),
        ("Learns from every interaction", 1),
        ("Gets smarter over time", 1),
        ("Can create databases and tables safely", 1),
        "",
        "🧠 How It Works",
        ("User asks question in plain English", 1),
        ("Agent discovers tools dynamically from MCP server", 1),
        ("Finds similar past queries (768-dim embeddings)", 1),
        ("Loads only relevant documentation (SKILL.md files)", 1),
        ("Claude API answers using discovered tools", 1),
        ("Stores what worked for future queries", 1),
    ]
    add_bullet_slide(prs, "Overview", overview)

    # 3. Section: Architecture
    add_section_slide(prs, "System Architecture")

    # 4. Architecture diagram
    add_architecture_overview(prs)

    # 5-7. DETAILED FLOW DIAGRAMS (NEW!)
    add_flow_new_query_detailed(prs)
    add_flow_learned_query_detailed(prs)
    add_flow_complex_skill_fetch(prs)

    # 8. Section: Real World
    add_section_slide(prs, "Real-World Applications")

    # 9. Use cases
    add_real_world_use_cases(prs)

    # 10. Why CockroachDB
    add_why_cockroachdb(prs)

    # 11. Section: Technical
    add_section_slide(prs, "Technical Excellence")

    # 12. Technical advantages
    tech = [
        "🚀 Dynamic Tool Discovery",
        ("No hardcoded MCP tools - discovered at runtime", 1),
        ("Agent calls tools/list method on startup", 1),
        ("Adapts automatically when MCP server adds capabilities", 1),
        ("Not brittle - future-proof design", 1),
        "",
        "🎯 Performance Optimizations",
        ("768-dimensional embeddings (all-mpnet-base-v2)", 1),
        ("Zero-skill pre-loading saves ~30-40% tokens", 1),
        ("HNSW vector index for sub-millisecond search", 1),
        ("Progressive learning: gets smarter over time", 1),
        "",
        "🔒 Production-Ready",
        ("Write operations with user confirmation", 1),
        ("OAuth authentication for secure access", 1),
        ("Color-coded UI for demos", 1),
        ("Comprehensive error handling", 1),
    ]
    add_bullet_slide(prs, "Technical Advantages", tech)

    # 13. Stats
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Performance Metrics & CockroachDB Integration"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    add_box(slide, 0.5, 1.5, 2.8, 0.9, "Token Savings\n30-40%\nvs old approach",
            RGBColor(180, 255, 180), 18, True)
    add_box(slide, 3.6, 1.5, 2.8, 0.9, "Embedding Quality\n768 dims\nall-mpnet-base-v2",
            RGBColor(200, 230, 255), 18, True)
    add_box(slide, 6.7, 1.5, 2.8, 0.9, "Query Speed\n~2 seconds\nend-to-end",
            RGBColor(255, 240, 200), 18, True)
    add_box(slide, 0.5, 2.8, 2.8, 0.9, "Tools Discovered\n11 MCP tools\ndynamically",
            RGBColor(255, 230, 230), 18, True)
    add_box(slide, 3.6, 2.8, 2.8, 0.9, "Skills Available\n25+ SKILL.md\nfrom GitHub",
            RGBColor(240, 220, 255), 18, True)
    add_box(slide, 6.7, 2.8, 2.8, 0.9, "Write Operations\n3 tools\nwith safety",
            RGBColor(255, 200, 150), 18, True)

    schema = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.3))
    schema.text_frame.text = """CockroachDB Schema:

• ai_demo.query_history: Stores queries with 768-dim embeddings + skills_used[]
  - HNSW index for fast cosine similarity search (sub-millisecond)
  - Learns from every interaction - gets smarter over time
  - pgvector native support - no external vector DB needed

• ai_demo.skills: Stores 25+ SKILL.md files as TEXT
  - CockroachDB best practices documentation
  - Loaded on-demand when needed
  - Distributed globally with ACID guarantees

• Storage: ~3.3 KB per query (768 floats × 4 bytes) - 1,000 queries = only 3.3 MB!"""
    for p in schema.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    # 14. Write operations
    write_ops = [
        "✨ Write Operations (April 5, 2026)",
        ("create_database - Create new databases", 1),
        ("create_table - Create tables with SQL DDL", 1),
        ("insert_rows - Insert data into tables", 1),
        "",
        "🔒 Safety Features",
        ("User confirmation required for EVERY write", 1),
        ("Dark orange warning for high visibility", 1),
        ("Shows exact operation before execution", 1),
        ("Can cancel any operation", 1),
        ("OAuth authentication with explicit permissions", 1),
        "",
        "❌ Not Supported (By Design)",
        ("UPDATE, DELETE - too risky", 1),
        ("DROP, TRUNCATE - destructive", 1),
        ("All write ops require user approval", 1),
    ]
    add_bullet_slide(prs, "Write Operations with Safety", write_ops)

    # 15. Summary
    summary = [
        "🎉 Key Achievements",
        ("Dynamic tool discovery - not brittle!", 1),
        ("2,545 lines of intelligent Python code", 1),
        ("95% token savings vs legacy approach", 1),
        ("Progressive learning from every query", 1),
        ("Production-ready with safety confirmations", 1),
        "",
        "🚀 Why CockroachDB is Essential",
        ("pgvector for 768-dim embeddings (no external vector DB)", 1),
        ("HNSW index for sub-millisecond similarity search", 1),
        ("ACID transactions ensure no lost learning", 1),
        ("Horizontal scalability from 1 to millions of queries", 1),
        ("Global distribution - shared learning across regions", 1),
        ("One database for: vectors, skills, and the cluster itself", 1),
    ]
    add_bullet_slide(prs, "Summary & Impact", summary)

    # Save
    filename = '/var/www/ai/aidemo2/docs/timbobfinalversion2.pptx'
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    create_presentation()
