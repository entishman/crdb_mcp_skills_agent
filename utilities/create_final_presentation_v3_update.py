#!/usr/bin/env python3
"""
Final Visual PowerPoint Presentation - Version 3 Update
Updated slides 6-8 with clearer flow progression
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def add_box(slide, x, y, width, height, text, fill_color, text_size=10, bold=False):
    """Add a colored box with BLACK text"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(2)

    box.text_frame.text = text
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = 1

    # Set ALL paragraphs to black (multi-line text creates multiple paragraphs)
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # BLACK TEXT!

    return box

def add_arrow(slide, x1, y1, x2, y2, label=None):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(2)
    connector.line.end_arrow_type = 2

    if label:
        label_box = slide.shapes.add_textbox(
            Inches((x1 + x2) / 2 - 0.4), Inches((y1 + y2) / 2 - 0.15),
            Inches(0.8), Inches(0.25)
        )
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.italic = True
        label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        label_box.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return connector

def add_slide_6_no_skill_needed(prs):
    """Slide 6: Simple task - Claude doesn't need SKILL.md"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 1: Simple Task - No SKILL.md Needed"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1: User
    box1 = add_box(slide, 0.5, 1, 2, 0.6, 'User Query\n"List all databases"',
            RGBColor(173, 216, 230), 11, True)

    # Step 2: MCP Discovery
    box2 = add_box(slide, 3, 1, 2, 0.6,
            'MCP Discovery\nmcp_client.py\nlist_tools()\n11 tools found',
            RGBColor(255, 210, 210), 9.5, False)
    add_arrow(slide, 2.5, 1.3, 3, 1.3)

    # Step 3: Check History
    box3 = add_box(slide, 5.5, 1, 2, 0.6,
            'Check History\nquery_store.py\nfind_similar_queries()\n❌ No match',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5, 1.3, 5.5, 1.3)

    # Step 4: Call Claude (no skills)
    box4 = add_box(slide, 7.75, 1, 2, 0.6,
            'Call Claude API\nagent.py\nNo skills pre-loaded\nJust tools available',
            RGBColor(200, 230, 255), 9.5, False)
    add_arrow(slide, 7.5, 1.3, 7.75, 1.3)

    # Step 5: Claude analyzes
    box5 = add_box(slide, 7.75, 2, 2, 0.7,
            'Claude Decides\n"This is simple"\n"I know list_databases"\n"No skill needed"',
            RGBColor(220, 255, 220), 9.5, True)
    add_arrow(slide, 8.75, 1.6, 8.75, 2)

    # Step 6: Execute MCP
    box6 = add_box(slide, 5.5, 2, 2, 0.7,
            'Execute MCP Tool\nmcp_client.py\nexecute_tool_call(\n"list_databases", {})',
            RGBColor(255, 210, 210), 9.5, False)
    add_arrow(slide, 7.75, 2.35, 7.5, 2.35)

    # Step 7: Results
    box7 = add_box(slide, 3, 2, 2, 0.7,
            'Database Results\ndefaultdb\nai_demo\npostgres\nsystem',
            RGBColor(180, 220, 180), 9.5, False)
    add_arrow(slide, 5.5, 2.35, 5, 2.35)

    # Step 8: Claude formats
    box8 = add_box(slide, 0.5, 2, 2, 0.7,
            'Claude Formats\n"Your cluster has 4\ndatabases:"\n(lists them)',
            RGBColor(255, 240, 200), 9.5, False)
    add_arrow(slide, 3, 2.35, 2.5, 2.35)

    # Step 9: Store Learning
    box9 = add_box(slide, 0.5, 3.1, 2, 0.65,
            'Store Query\nquery_store.py\nstore_query()\nskills_used: [] ✅',
            RGBColor(220, 255, 220), 9.5, False)
    add_arrow(slide, 1.5, 2.7, 1.5, 3.1)

    # Key insight box
    insight = slide.shapes.add_textbox(Inches(3), Inches(3.1), Inches(6.75), Inches(0.65))
    insight.text_frame.text = "🎯 Key Point: Claude uses its training + tool descriptions to handle simple tasks.\nNo SKILL.md needed! Stored with empty skills_used[] for future reference."
    for p in insight.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Flow summary
    summary = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.3))
    summary.text_frame.text = """Complete Flow Summary:

1. User asks simple question → 2. MCP discovers 11 tools → 3. Check history (no match)
4. Call Claude with tools only (no skills) → 5. Claude: "I know this, use list_databases"
6. Execute MCP tool → 7. Get results from CockroachDB → 8. Claude formats answer
9. Store: query + embedding + skills_used:[] (empty - no skills needed!)

📊 Stored in ai_demo.query_history:
   - query_text: "List all databases"
   - embedding: [768-dimensional vector]
   - skills_used: []  ← Empty because Claude didn't need help!
   - timestamp: 2026-04-05 10:23:15

💡 Next time similar query: Still no skills loaded (because none were used)
   Simple queries stay fast and lightweight!"""
    for p in summary.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_slide_7_claude_fetches_skill(prs):
    """Slide 7: Complex task - Claude calls fetch_skill, stores result"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 2: Complex Task - Claude Fetches SKILL.md (No History Match)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    box1 = add_box(slide, 0.4, 0.7, 2, 0.6, 'User Query\n"Set up MOLT\nReplicator"',
            RGBColor(173, 216, 230), 11, True)

    # Step 2
    box2 = add_box(slide, 2.7, 0.7, 1.8, 0.6,
            'MCP Discovery\n11 tools\nincluding\nfetch_skill',
            RGBColor(255, 210, 210), 9, False)
    add_arrow(slide, 2.4, 1, 2.7, 1)

    # Step 3
    box3 = add_box(slide, 4.8, 0.7, 1.8, 0.6,
            'Check History\nfind_similar()\n❌ No match\n(never asked)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 4.5, 1, 4.8, 1)

    # Step 4
    box4 = add_box(slide, 7, 0.7, 2.6, 0.6,
            'Call Claude\nNo skills pre-loaded\nTools: MCP + fetch_skill\nUser query included',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 6.6, 1, 7, 1)

    # Step 5
    box5 = add_box(slide, 7, 1.7, 2.6, 0.7,
            'Claude Analyzes\n"MOLT migration is complex"\n"I need official guide"\nDecides: call fetch_skill',
            RGBColor(255, 240, 150), 9, True)
    add_arrow(slide, 8.3, 1.3, 8.3, 1.7)

    # Step 6
    box6 = add_box(slide, 4.8, 1.7, 1.8, 0.7,
            'Claude Calls\nfetch_skill(\n"onboarding-\nand-migrations/\nmolt-replicator")',
            RGBColor(255, 200, 150), 8.5, True)
    add_arrow(slide, 7, 2.05, 6.6, 2.05)

    # Step 7
    box7 = add_box(slide, 2.7, 1.7, 1.8, 0.7,
            'Fetch SKILL.md\nskill_fetcher.py\nQuery ai_demo\n.skills table\nReturns 8KB',
            RGBColor(255, 240, 200), 8.5, False)
    add_arrow(slide, 4.8, 2.05, 4.5, 2.05)

    # Step 8
    box8 = add_box(slide, 0.4, 1.7, 2, 0.7,
            'SKILL.md Loaded\nBest practices\nSetup steps\nCommon issues',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 2.7, 2.05, 2.4, 2.05)

    # Step 9
    box9 = add_box(slide, 0.4, 2.8, 2, 0.65,
            'Re-process\nClaude now has:\n• Tools\n• SKILL.md\n• User query',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 1.4, 2.4, 1.4, 2.8)

    # Step 10
    box10 = add_box(slide, 2.7, 2.8, 1.8, 0.65,
            'Expert Answer\nClaude uses\nSKILL.md to\nguide setup',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.4, 3.1, 2.7, 3.1)

    # Step 11
    box11 = add_box(slide, 4.8, 2.8, 2.2, 0.65,
            'Store Learning ✅\nquery + embedding +\nskills_used: ["molt-replicator"]\nStored in query_history',
            RGBColor(180, 255, 180), 8.5, True)
    add_arrow(slide, 4.5, 3.1, 4.8, 3.1)

    # Key insight
    insight = slide.shapes.add_textbox(Inches(0.4), Inches(3.7), Inches(9.2), Inches(0.6))
    insight.text_frame.text = '🎯 Critical Learning Moment: Query stored with skills_used:["onboarding-and-migrations/molt-replicator"]\nNext similar query will pre-load this skill automatically!'
    for p in insight.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 0, 102)

    # Detailed storage
    storage = slide.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(9.2), Inches(1.9))
    storage.text_frame.text = """What Gets Stored in ai_demo.query_history:

📊 Database Row Created:
   - query_text: "Set up MOLT Replicator"
   - embedding: [768-dim vector from sentence-transformers]
   - skills_used: ["onboarding-and-migrations/molt-replicator"]  ← THIS IS KEY!
   - timestamp: 2026-04-05 11:45:22

🔍 Why This Matters:
   • Future queries like "help with MOLT setup" or "MOLT migration guide"
   • pgvector similarity search finds this stored query (cosine similarity > 0.3)
   • Agent sees: "Oh, someone asked about MOLT before and used molt-replicator skill"
   • Pre-loads that skill BEFORE calling Claude → Saves 1 API call!

💡 Claude made the decision: "I need molt-replicator skill"
   Agent just executed and remembered it for next time."""
    for p in storage.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_slide_8_preloaded_from_history(prs):
    """Slide 8: Similar query - skill pre-loaded from history"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 3: Optimization - Skill Pre-loaded from Query History"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    box1 = add_box(slide, 0.4, 0.7, 2, 0.6, 'User Query\n"Help with MOLT\nmigration"',
            RGBColor(173, 216, 230), 11, True)

    # Step 2
    box2 = add_box(slide, 2.7, 0.7, 1.8, 0.6,
            'Generate\nEmbedding\n768-dim\nvector',
            RGBColor(255, 230, 230), 9, False)
    add_arrow(slide, 2.4, 1, 2.7, 1)

    # Step 3 - MATCH FOUND!
    box3 = add_box(slide, 4.8, 0.7, 2.2, 0.6,
            'Check History\npgvector search\n✅ MATCH FOUND!\nSimilarity: 0.82',
            RGBColor(180, 255, 180), 9, True)
    add_arrow(slide, 4.5, 1, 4.8, 1)

    # Step 4
    box4 = add_box(slide, 7.3, 0.7, 2.3, 0.6,
            'Found Match:\n"Set up MOLT Replicator"\nskills_used:\n["molt-replicator"]',
            RGBColor(220, 255, 220), 8.5, False)
    add_arrow(slide, 7, 1, 7.3, 1)

    # Step 5
    box5 = add_box(slide, 7.3, 1.7, 2.3, 0.7,
            'Pre-load Skill\nskill_fetcher.py\nget_skill(\n"molt-replicator")\nBEFORE Claude call',
            RGBColor(255, 240, 200), 8.5, True)
    add_arrow(slide, 8.45, 1.3, 8.45, 1.7)

    # Step 6
    box6 = add_box(slide, 4.8, 1.7, 2.2, 0.7,
            'Call Claude\nWith skill\nPRE-LOADED\n+ tools\n+ query',
            RGBColor(200, 230, 255), 9, True)
    add_arrow(slide, 7.3, 2.05, 7, 2.05)

    # Step 7
    box7 = add_box(slide, 2.7, 1.7, 1.8, 0.7,
            'Claude Has\nEverything!\nNo need to\ncall fetch_skill',
            RGBColor(220, 255, 220), 9, True)
    add_arrow(slide, 4.8, 2.05, 4.5, 2.05)

    # Step 8
    box8 = add_box(slide, 0.4, 1.7, 2, 0.7,
            'Expert Answer\nImmediately\nUsing SKILL.md\nguidance',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.7, 2.05, 2.4, 2.05)

    # Step 9
    box9 = add_box(slide, 0.4, 2.8, 2, 0.65,
            'Store Query\nConfirm skill\nhelped again\nReinforces',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 1.4, 2.4, 1.4, 2.8)

    # Performance comparison
    perf = slide.shapes.add_textbox(Inches(2.7), Inches(2.8), Inches(6.9), Inches(0.65))
    perf.text_frame.text = "⚡ Performance: 1 Claude API call (vs 2 in Flow 2)\n💰 Cost Savings: ~40% fewer tokens | ⏱️ Time Savings: ~1 second faster"
    for p in perf.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Detailed comparison
    comparison = slide.shapes.add_textbox(Inches(0.4), Inches(3.7), Inches(9.2), Inches(2.7))
    comparison.text_frame.text = """🔍 What Happened Behind the Scenes:

Step 3 Details - pgvector Similarity Search:
   • Query "Help with MOLT migration" → embedding generated
   • CockroachDB executes: SELECT query_text, skills_used, embedding <-> [new_vector] AS similarity
                           FROM ai_demo.query_history
                           WHERE embedding <-> [new_vector] < 0.7
                           ORDER BY similarity LIMIT 5
   • HNSW index makes this sub-millisecond!
   • Best match: "Set up MOLT Replicator" (similarity 0.82)
   • Extract skills_used: ["onboarding-and-migrations/molt-replicator"]

Step 5 - Optimization in Action:
   • Agent: "I'll load molt-replicator skill BEFORE calling Claude"
   • Fetches 8KB SKILL.md from ai_demo.skills table
   • Includes it in Claude API call's system prompt

The Result:
   ❌ Flow 2 (no history): User → Claude → fetch_skill → Claude again → Answer (2 API calls)
   ✅ Flow 3 (with history): User → [pre-load skill] → Claude → Answer (1 API call)

💡 This is the Learning Optimization: Same quality answer, half the API calls!"""
    for p in comparison.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide


if __name__ == "__main__":
    """Load existing presentation and update slides 6-8"""
    import sys

    # Load the existing presentation
    input_file = '/var/www/ai/aidemo2/docs/timbobfinalversion3.pptx'
    output_file = '/var/www/ai/aidemo2/docs/timbobfinalversion3-update.pptx'

    try:
        prs = Presentation(input_file)
        print(f"✅ Loaded existing presentation: {input_file}")
        print(f"📊 Current slide count: {len(prs.slides)}")

        # Delete slides 5, 6, 7 (indexes 5, 6, 7 since 0-indexed)
        # We need to delete slide 7 first, then 6, then 5 (reverse order)
        rId = prs.slides._sldIdLst[7].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[7]

        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

        rId = prs.slides._sldIdLst[5].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[5]

        print(f"✅ Deleted old slides 6-8")
        print(f"📊 Slide count after deletion: {len(prs.slides)}")

        # Now we need to insert new slides at position 5 (after slide 4, the architecture section)
        # Unfortunately python-pptx doesn't support inserting slides at arbitrary positions easily
        # So we'll need to add them at the end and then move them
        # Actually, let's just rebuild from the template

        print("⚠️ Note: python-pptx doesn't support easy slide insertion.")
        print("Creating new presentation with updated slides 6-8...")

        # Let's create fresh slides
        add_slide_6_no_skill_needed(prs)
        add_slide_7_claude_fetches_skill(prs)
        add_slide_8_preloaded_from_history(prs)

        # Save
        prs.save(output_file)
        print(f"✅ Presentation created: {output_file}")
        print(f"📊 Total slides: {len(prs.slides)}")
        print()
        print("⚠️ NOTE: New slides were added at the END.")
        print("You may need to manually reorder them in PowerPoint.")
        print("Or I can regenerate the entire presentation with proper ordering.")

    except Exception as e:
        print(f"❌ Error: {e}")
        print()
        print("I'll create a standalone presentation with just the 3 updated slides for review:")

        # Create fresh presentation with just these 3 slides
        prs_new = Presentation()
        prs_new.slide_width = Inches(10)
        prs_new.slide_height = Inches(7.5)

        add_slide_6_no_skill_needed(prs_new)
        add_slide_7_claude_fetches_skill(prs_new)
        add_slide_8_preloaded_from_history(prs_new)

        review_file = '/var/www/ai/aidemo2/docs/slides_6-8_review.pptx'
        prs_new.save(review_file)
        print(f"✅ Review slides created: {review_file}")
        print(f"📊 Contains 3 slides for review")
