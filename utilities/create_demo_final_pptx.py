#!/usr/bin/env python3
"""
Create demo_final.pptx - Complete presentation with all updates

This script creates a comprehensive PowerPoint presentation combining:
- Content from existing demo_for_humanx.pptx
- New write operations features
- All relevant content from documentation files
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx library not installed")
    print("Install with: pip install python-pptx")
    print("\nAlternatively, use the detailed content in:")
    print("  - POWERPOINT_SLIDES_READY_TO_ADD.md")
    print("  - PRESENTATION_UPDATES.md")
    exit(1)

import os
from pathlib import Path


def create_presentation():
    """Create the complete demo_final.pptx presentation"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide = add_title_slide(prs,
        "CockroachDB Intelligent Agent",
        "AI-Powered Database Management with Write Operations\nApril 5, 2026"
    )

    # Slide 2: What's New Today
    slide = add_bullet_slide(prs,
        "What's New Today! 🎉",
        [
            "✅ Write Operations Enabled - Create databases and tables",
            "✅ OAuth Authentication - Secure CockroachDB Cloud integration",
            "✅ Safety Confirmations - User approval required for all writes",
            "✅ 20 Unique Skill Tests - Comprehensive verification",
            "✅ Enhanced Documentation - Complete implementation guides"
        ]
    )

    # Slide 3: Agent Capabilities
    slide = add_bullet_slide(prs,
        "What This Agent Can Do",
        [
            "✅ Answers CockroachDB questions in natural language",
            "✅ Learns from every interaction (gets smarter over time)",
            "✅ Uses semantic embeddings to find similar queries",
            "✅ Loads only needed documentation (95% token savings)",
            "✅ NEW: Create databases with user confirmation",
            "✅ NEW: Create tables with schema validation",
            "✅ NEW: Insert rows with type checking"
        ]
    )

    # Slide 4: Write Operations Enabled
    slide = add_bullet_slide(prs,
        "Write Operations - Now Available",
        [
            "✅ Create Databases - Natural language database creation",
            "✅ Create Tables - Schema design with AI assistance",
            "✅ Insert Rows - Data insertion via conversational interface",
            "✅ Safety First - User confirmation required before execution",
            "✅ OAuth Authentication - Secure access via CockroachDB Cloud",
            "",
            "❌ No destructive operations (DROP, DELETE, UPDATE)",
            "❌ No cluster configuration changes"
        ]
    )

    # Slide 5: Multi-Layer Safety
    slide = add_bullet_slide(prs,
        "Multi-Layer Safety Architecture",
        [
            "Layer 1: OAuth 2.1 Authentication",
            "  • User must log in and grant permissions explicitly",
            "  • Short-lived tokens (more secure than API keys)",
            "",
            "Layer 2: Interactive Confirmation",
            "  • Every write operation requires yes/no confirmation",
            "  • Dark orange warning for high visibility",
            "",
            "Layer 3: Limited Scope",
            "  • Only CREATE DATABASE, CREATE TABLE, INSERT ROWS",
            "  • No destructive operations allowed",
            "",
            "Layer 4: Role-Based Access",
            "  • Requires Cluster Admin or Cluster Operator role"
        ]
    )

    # Slide 6: How It Works
    slide = add_bullet_slide(prs,
        "From Request to Database Creation",
        [
            "1. User Request",
            "   \"Create a database called production\"",
            "",
            "2. Claude API analyzes intent",
            "   Recognizes create_database operation needed",
            "",
            "3. ⚠️  WRITE OPERATION REQUESTED (dark orange)",
            "   Tool: create_database",
            "   Input: {\"database\": \"production\"}",
            "",
            "4. User confirms: yes",
            "",
            "5. ✓ Database created successfully!"
        ]
    )

    # Slide 7: Live Demo
    slide = add_bullet_slide(prs,
        "Live Demo - Creating a Database",
        [
            "1. Start agent: python3 agent.py",
            "",
            "2. Request: create a database called timbob",
            "",
            "3. Confirmation prompt appears (dark orange):",
            "   ⚠️  WRITE OPERATION REQUESTED",
            "   Execute this operation? (yes/no):",
            "",
            "4. User types: yes",
            "",
            "5. ✓ Write operation completed",
            "   Database 'timbob' created successfully!"
        ]
    )

    # Slide 8: How the Agent Learns
    slide = add_bullet_slide(prs,
        "Intelligent Learning System",
        [
            "Query-Based Learning:",
            "  • Remembers which skills helped with each query",
            "  • Pre-loads proven skills for similar questions",
            "  • Gets smarter over time",
            "",
            "Semantic Understanding:",
            "  • 768-dimensional embeddings (all-mpnet-base-v2)",
            "  • Understands paraphrases and intent",
            "  • Similarity search with HNSW index",
            "",
            "Zero-Skill Optimization:",
            "  • Simple queries: 0 skills loaded",
            "  • Complex queries: Fetch skills on-demand",
            "  • 30-40% token savings"
        ]
    )

    # Slide 9: Technical Architecture
    slide = add_bullet_slide(prs,
        "System Architecture",
        [
            "Frontend: Natural language interface",
            "  ↓",
            "Claude API (Sonnet 4.6): Intent understanding",
            "  ↓",
            "Agent Logic: Query learning + Skill fetching",
            "  ↓",
            "MCP Client: Tool execution",
            "  ↓",
            "CockroachDB Cloud: Database operations",
            "",
            "Key Components:",
            "• sentence-transformers (768-dim embeddings)",
            "• CockroachDB pgvector (similarity search)",
            "• OAuth 2.1 (secure authentication)",
            "• Model Context Protocol (tool integration)"
        ]
    )

    # Slide 10: Performance Metrics
    slide = add_bullet_slide(prs,
        "Performance & Efficiency",
        [
            "Token Efficiency:",
            "  • 95% savings vs loading all docs",
            "  • Additional 30-40% savings from zero-skill optimization",
            "",
            "Response Time:",
            "  • ~2 seconds average",
            "  • Includes embedding generation and similarity search",
            "",
            "Learning:",
            "  • Automatic from every query",
            "  • No manual training required",
            "",
            "Code Size:",
            "  • 2,545 lines of intelligent Python",
            "  • Fully documented and tested"
        ]
    )

    # Slide 11: Technical Implementation
    slide = add_bullet_slide(prs,
        "What Changed Under the Hood",
        [
            "BEFORE (Read-Only):",
            "  • System prompt: 'CANNOT execute write operations'",
            "  • Tool list: get_all_tools() (read-only only)",
            "  • No confirmation logic",
            "",
            "AFTER (Write-Enabled):",
            "  • System prompt: 'You have access to write operations'",
            "  • Tool list: get_tool_definitions() (all tools)",
            "  • Safety confirmation for all writes",
            "  • OAuth authentication",
            "  • Dark orange warnings for visibility",
            "",
            "Code Changes:",
            "  • agent.py line 700: Tool definitions expanded",
            "  • agent.py lines 785-801: Confirmation logic added",
            "  • agent.py lines 306-325: System prompt updated"
        ]
    )

    # Slide 12: Documentation
    slide = add_bullet_slide(prs,
        "Complete Documentation",
        [
            "Setup Guides:",
            "  • COCKROACHDB_SETUP.md (25KB) - OAuth + MCP setup",
            "  • REQUIREMENTS_WRITE_OPERATIONS.md (32KB) - Implementation",
            "",
            "Testing:",
            "  • SKILL_TEST_TASKS.md - 20 unique verification tasks",
            "  • MANUAL_SKILL_TEST_GUIDE.md - Testing instructions",
            "",
            "Architecture:",
            "  • ARCHITECTURE.md - Complete system design",
            "  • CHANGELOG.md - All changes documented",
            "",
            "Code Documentation:",
            "  • COMMENTS_agent.md - Line-by-line explanation",
            "  • Inline comments throughout codebase"
        ]
    )

    # Slide 13: Skill Testing
    slide = add_bullet_slide(prs,
        "20 Unique Skills Verified",
        [
            "Observability & Diagnostics (7 skills):",
            "  • Range distribution, schema changes, table stats",
            "  • Background jobs, statement/transaction profiling",
            "  • Live SQL activity triaging",
            "",
            "Operations & Lifecycle (7 skills):",
            "  • Certificates, capacity, settings, maintenance",
            "  • Production provisioning, health review, upgrades",
            "",
            "Onboarding & Migrations (3 skills):",
            "  • MOLT Fetch, MOLT Replicator, MOLT Verify",
            "",
            "Security & Governance (2 skills):",
            "  • Security audit, audit logging",
            "",
            "Query & Schema Design (1 skill):",
            "  • CockroachDB SQL best practices"
        ]
    )

    # Slide 14: What's Next
    slide = add_bullet_slide(prs,
        "Future Enhancements",
        [
            "Potential Additions:",
            "  • UPDATE and DELETE operations (with enhanced safety)",
            "  • Schema migrations (ALTER TABLE)",
            "  • Batch operations",
            "  • Transaction support",
            "  • Rollback capabilities",
            "  • Operation history and audit log",
            "",
            "Current Focus:",
            "  ✅ Write operations with safety (COMPLETE)",
            "  ✅ Query learning optimization (COMPLETE)",
            "  ✅ Vector search improvements (COMPLETE)",
            "  • Performance monitoring dashboard",
            "  • Multi-cluster support"
        ]
    )

    # Slide 15: How to Use
    slide = add_bullet_slide(prs,
        "Getting Started",
        [
            "For Read Operations:",
            "  python3 agent.py",
            "  • List databases",
            "  • Show tables",
            "  • Execute SELECT queries",
            "  • Review cluster health",
            "",
            "For Write Operations:",
            "  python3 agent.py",
            "  • 'Create a database called myapp'",
            "  • 'Create a users table with id and email'",
            "  • 'Insert a test user'",
            "  • Review operation → Confirm → Execute",
            "",
            "Requirements:",
            "  • CockroachDB Cloud cluster",
            "  • OAuth authentication configured",
            "  • Cluster Admin or Cluster Operator role"
        ]
    )

    # Slide 16: Key Achievements
    slide = add_bullet_slide(prs,
        "Key Achievements",
        [
            "✅ Self-learning AI that gets smarter over time",
            "✅ 95% token savings through dynamic skill loading",
            "✅ 768-dimensional semantic embeddings",
            "✅ Zero-skill optimization (30-40% additional savings)",
            "✅ Production-safe write operations",
            "✅ Multi-layer security architecture",
            "✅ Complete documentation (50KB+)",
            "✅ 20 unique skills tested and verified",
            "✅ Color-coded UI for demos",
            "✅ OAuth authentication integrated",
            "",
            "2,545 lines of intelligent, documented code",
            "Ready for production use"
        ]
    )

    # Final slide: Q&A
    slide = add_title_slide(prs,
        "Questions?",
        "CockroachDB Intelligent Agent\nDemo & Discussion"
    )

    return prs


def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide


def add_bullet_slide(prs, title, bullets):
    """Add a slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0 if not bullet.startswith('  ') else 1
        p.font.size = Pt(18) if p.level == 0 else Pt(16)

    return slide


def main():
    """Main function to create the presentation"""
    print("=" * 80)
    print("Creating demo_final.pptx")
    print("=" * 80)
    print()

    # Create presentation
    print("Building presentation with all content...")
    prs = create_presentation()

    # Save presentation
    output_path = Path(__file__).parent.parent / "docs" / "demo_final.pptx"
    print(f"Saving to: {output_path}")
    prs.save(str(output_path))

    print()
    print("=" * 80)
    print("✓ SUCCESS: demo_final.pptx created!")
    print("=" * 80)
    print()
    print(f"Location: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    print()
    print("Content includes:")
    print("  • All updated capabilities")
    print("  • Write operations features")
    print("  • Safety architecture")
    print("  • Live demo script")
    print("  • Technical implementation")
    print("  • Complete documentation")
    print("  • Skill testing results")
    print("  • Future roadmap")
    print()
    print("Ready to present! 🎉")
    print()


if __name__ == "__main__":
    main()
