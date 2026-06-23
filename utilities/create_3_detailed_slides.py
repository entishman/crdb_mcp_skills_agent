#!/usr/bin/env python3
"""
Create 3 detailed technical slides with function names and architecture details
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_box(slide, x, y, width, height, text, fill_color, text_size=10, bold=False):
    """Add a colored box with black text"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(2)

    box.text_frame.text = text
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = 1

    # Set ALL paragraphs to black (multi-line text creates multiple paragraphs)
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # BLACK TEXT!

    return box

def add_arrow(slide, x1, y1, x2, y2, label=None):
    """Add an arrow connector"""
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(2)
    connector.line.end_arrow_type = 2

    if label:
        label_box = slide.shapes.add_textbox(
            Inches((x1 + x2) / 2 - 0.4), Inches((y1 + y2) / 2 - 0.15),
            Inches(0.8), Inches(0.25)
        )
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.italic = True
        label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        label_box.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return connector

def create_slide_1_new_query_detailed(prs):
    """Slide 1: New Query (Cold Start) - Detailed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 1: New Query (Cold Start) - No Learning Yet"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1: User query
    add_box(slide, 0.3, 0.7, 2.2, 0.55,
            'Step 1: User Input\n"list databases"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    # Step 2: Generate embedding - DETAILED
    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'ollama_embedding_manager.py\n' +
            'generate_embedding()\n' +
            '→ 768-dim vector',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3: Search query_history - DETAILED
    add_box(slide, 5.6, 0.7, 2.9, 0.55,
            'Step 3: Search Similar Queries\n' +
            'query_store.py: find_similar_queries()\n' +
            'DB: CockroachDB ai_demo.query_history\n' +
            'Result: ❌ No match (new query)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 7, 1.25, 7, 1.5)

    # Step 4: Pre-load skills - DETAILED
    add_box(slide, 5.6, 1.5, 2.9, 0.55,
            'Step 4: Pre-load Skills\n' +
            'agent.py: _build_system_prompt()\n' +
            'calls _get_default_skills()\n' +
            'Returns: [] (zero skills!)',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 7, 2.05, 7, 2.3)

    # Step 5: Build prompt & call Claude - DETAILED
    add_box(slide, 5.6, 2.3, 2.9, 0.65,
            'Step 5: Call Claude API\n' +
            'agent.py: run_task()\n' +
            'claude.messages.create(\n' +
            '  system=prompt,\n' +
            '  tools=get_all_tools()\n' +
            ')',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 7, 2.95, 7, 3.2)

    # Step 6: Claude uses MCP tools - DETAILED
    add_box(slide, 5.6, 3.2, 2.9, 0.75,
            'Step 6: Claude Selects Tool\n' +
            'Tool definitions sent in Step 5\n' +
            'Claude sees: list_databases tool\n' +
            'Description: "Lists all databases"\n' +
            'Claude calls: list_databases()',
            RGBColor(240, 220, 255), 8.5, False)
    add_arrow(slide, 7, 3.95, 7, 4.2)

    # Step 7: MCP executes query - DETAILED
    add_box(slide, 5.6, 4.2, 2.9, 0.65,
            'Step 7: Execute via MCP\n' +
            'agent.py: calls mcp_client.py\n' +
            'execute_tool_call("list_databases")\n' +
            'MCP Server → CockroachDB\n' +
            'Runs: SHOW DATABASES',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 5.6, 4.5, 5.3, 4.5)

    # Step 8: Return results
    add_box(slide, 2.8, 4.2, 2.5, 0.65,
            'Step 8: Results to Claude\n' +
            'MCP returns database list\n' +
            'Claude formats answer',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.8, 4.5, 2.5, 4.5)

    # Step 9: Store learning
    add_box(slide, 0.3, 4.2, 2.2, 0.65,
            'Step 9: Store Learning\n' +
            'query_store.py: store_query()\n' +
            'Saves: query + [] skills\n' +
            'DB: ai_demo.query_history',
            RGBColor(220, 255, 220), 9, False)

    # Add detailed explanation boxes
    explain1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.1), Inches(4.5), Inches(0.6))
    explain1.text_frame.text = "💡 MCP Tool Discovery: When Claude connects, it receives tool definitions from get_all_tools() which describes each tool's name, description, and parameters. Example: {'name': 'list_databases', 'description': 'Lists all databases in the cluster', 'input_schema': {...}}"
    explain1.text_frame.paragraphs[0].font.size = Pt(8.5)
    explain1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    explain1.text_frame.word_wrap = True

    explain2 = slide.shapes.add_textbox(Inches(5), Inches(5.1), Inches(4.5), Inches(0.6))
    explain2.text_frame.text = "🎯 Zero-Skill Optimization: Simple queries like 'list databases' don't need SKILL.md files. Claude knows basic SQL from training. This saves ~8KB tokens per query. Only complex operations trigger skill fetching."
    explain2.text_frame.paragraphs[0].font.size = Pt(8.5)
    explain2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    explain2.text_frame.word_wrap = True

    # Add legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(5.85), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "Database Location: ai_demo.query_history and ai_demo.skills tables are in the same CockroachDB cluster being queried"
    legend.text_frame.paragraphs[0].font.size = Pt(9)
    legend.text_frame.paragraphs[0].font.italic = True
    legend.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

    return slide

def create_slide_2_learned_query_detailed(prs):
    """Slide 2: Learned Query (Warm Start) - Detailed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 2: Learned Query (Warm Start) - Learning Kicks In"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1: User query
    add_box(slide, 0.3, 0.7, 2.2, 0.55,
            'Step 1: User Input\n"check for hotspots"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    # Step 2: Generate embedding
    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'ollama_embedding_manager.py\n' +
            'generate_embedding()\n' +
            '→ 768-dim vector',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3: Search query_history - MATCH FOUND!
    add_box(slide, 5.6, 0.7, 2.9, 0.6,
            'Step 3: Search Similar Queries\n' +
            'query_store.py: find_similar_queries()\n' +
            'SQL: cosine similarity search\n' +
            'Result: ✅ Match! (0.78 similarity)\n' +
            'Matched: "analyze hot ranges"',
            RGBColor(180, 255, 180), 8.5, False)
    add_arrow(slide, 7, 1.3, 7, 1.55)

    # Step 4: Load learned skills
    add_box(slide, 5.6, 1.55, 2.9, 0.6,
            'Step 4: Load Learned Skills\n' +
            'agent.py: _build_system_prompt()\n' +
            'Gets skills from matched query:\n' +
            '["observability-and-diagnostics/\n' +
            'analyzing-range-distribution"]',
            RGBColor(220, 255, 220), 8.5, False)
    add_arrow(slide, 7, 2.15, 7, 2.4)

    # Step 4b: Fetch skill content
    add_box(slide, 5.6, 2.4, 2.9, 0.6,
            'Step 4b: Fetch SKILL.md\n' +
            'skill_fetcher.py: get_skill()\n' +
            'Source: ai_demo.skills table\n' +
            'Returns: Full SKILL.md content\n' +
            '(Best practices for range analysis)',
            RGBColor(255, 240, 200), 8.5, False)
    add_arrow(slide, 7, 3, 7, 3.25)

    # Step 5: Call Claude with skill context
    add_box(slide, 5.6, 3.25, 2.9, 0.65,
            'Step 5: Call Claude API\n' +
            'agent.py: run_task()\n' +
            'System prompt includes:\n' +
            '  - SKILL.md content\n' +
            '  - Tool definitions\n' +
            '  - User query',
            RGBColor(200, 230, 255), 8.5, False)
    add_arrow(slide, 7, 3.9, 7, 4.15)

    # Step 6: Claude uses skill knowledge
    add_box(slide, 5.6, 4.15, 2.9, 0.6,
            'Step 6: Claude Uses Skill\n' +
            'References SKILL.md best practices\n' +
            'Knows to check: SHOW RANGES\n' +
            'Calls: select_query() with\n' +
            'range distribution query',
            RGBColor(240, 220, 255), 8.5, False)
    add_arrow(slide, 5.6, 4.45, 5.3, 4.45)

    # Step 7: Execute query
    add_box(slide, 2.8, 4.15, 2.5, 0.6,
            'Step 7: Execute Query\n' +
            'mcp_client.py: execute_tool_call()\n' +
            'MCP Server queries cluster\n' +
            'Returns: range distribution data',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 2.8, 4.45, 2.5, 4.45)

    # Step 8: Expert answer
    add_box(slide, 0.3, 4.15, 2.2, 0.6,
            'Step 8: Expert Answer\n' +
            'Claude analyzes results\n' +
            'using SKILL.md guidance\n' +
            'Provides detailed analysis',
            RGBColor(255, 240, 200), 9, False)

    # Step 9: Update learning
    add_box(slide, 0.3, 5, 2.2, 0.5,
            'Step 9: Confirm Learning\n' +
            'query_store.py: store_query()\n' +
            'Confirms skill helped',
            RGBColor(220, 255, 220), 9, False)

    # Explanation
    explain = slide.shapes.add_textbox(Inches(2.8), Inches(5), Inches(6.7), Inches(0.5))
    explain.text_frame.text = "🎯 Progressive Learning: Each query improves future performance. Similarity threshold 0.3 catches paraphrases. After 50+ queries, most questions have learned matches = instant expert answers!"
    explain.text_frame.paragraphs[0].font.size = Pt(9)
    explain.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    explain.text_frame.word_wrap = True

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(5.7), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "Similarity Search: Uses pgvector cosine distance: 1 - (embedding <-> query_embedding). HNSW index makes search sub-millisecond even with 1M+ queries"
    legend.text_frame.paragraphs[0].font.size = Pt(9)
    legend.text_frame.paragraphs[0].font.italic = True
    legend.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

    return slide

def create_slide_3_complex_skill_fetch(prs):
    """Slide 3: Complex Query Requiring Skill Fetch - Detailed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 3: Complex Query (No Match) - Dynamic Skill Fetching"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1: User query - complex
    add_box(slide, 0.3, 0.7, 2.2, 0.55,
            'Step 1: User Input\n"How do I set up\nMOLT Replicator?"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    # Step 2: Generate embedding
    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'ollama_embedding_manager.py\n' +
            'generate_embedding()\n' +
            '→ 768-dim vector',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3: Search - no match
    add_box(slide, 5.6, 0.7, 2.9, 0.55,
            'Step 3: Search Query History\n' +
            'query_store.py: find_similar_queries()\n' +
            'Result: ❌ No match\n' +
            '(Never asked about MOLT before)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 7, 1.25, 7, 1.5)

    # Step 4: Zero skills pre-loaded
    add_box(slide, 5.6, 1.5, 2.9, 0.5,
            'Step 4: Pre-load Skills\n' +
            'agent.py: _get_default_skills()\n' +
            'Returns: [] (zero skills)',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 7, 2, 7, 2.25)

    # Step 5: Call Claude - first attempt
    add_box(slide, 5.6, 2.25, 2.9, 0.6,
            'Step 5: Call Claude API\n' +
            'agent.py: run_task()\n' +
            'Claude gets: query + tool list\n' +
            'Claude thinks: "I need more info\n' +
            'about MOLT Replicator"',
            RGBColor(200, 230, 255), 8.5, False)
    add_arrow(slide, 7, 2.85, 7, 3.1)

    # Step 6: Claude calls fetch_skill
    add_box(slide, 5.6, 3.1, 2.9, 0.6,
            'Step 6: Claude Calls fetch_skill\n' +
            'Tool use: fetch_skill(\n' +
            '  skill_name="onboarding-and-\n' +
            '  migrations/molt-replicator"\n' +
            ')',
            RGBColor(255, 200, 150), 8.5, False)
    add_arrow(slide, 5.6, 3.4, 5.3, 3.4)

    # Step 7: Agent fetches skill
    add_box(slide, 2.8, 3.1, 2.5, 0.6,
            'Step 7: Fetch SKILL.md\n' +
            'agent.py: _handle_fetch_skill()\n' +
            'skill_fetcher.py: get_skill()\n' +
            'Source: ai_demo.skills table\n' +
            'Returns: Full SKILL.md content',
            RGBColor(255, 240, 200), 8, False)
    add_arrow(slide, 2.8, 3.4, 2.5, 3.4)

    # Step 8: Return skill to Claude
    add_box(slide, 0.3, 3.1, 2.2, 0.6,
            'Step 8: Skill → Claude\n' +
            'agent.py appends skill\n' +
            'to conversation\n' +
            'Claude now has expertise!',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 1.4, 3.7, 1.4, 3.95)

    # Step 9: Claude uses skill
    add_box(slide, 0.3, 3.95, 2.2, 0.6,
            'Step 9: Claude Re-processes\n' +
            'Now with SKILL.md context\n' +
            'Knows MOLT setup steps\n' +
            'May call MCP tools if needed',
            RGBColor(240, 220, 255), 9, False)
    add_arrow(slide, 2.5, 4.25, 2.8, 4.25)

    # Step 10: Execute operations
    add_box(slide, 2.8, 3.95, 2.5, 0.6,
            'Step 10: Execute if Needed\n' +
            'May check existing setup:\n' +
            'mcp_client.py: execute_tool_call()\n' +
            'list_databases(), etc.',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 5.3, 4.25, 5.6, 4.25)

    # Step 11: Expert answer
    add_box(slide, 5.6, 3.95, 2.9, 0.6,
            'Step 11: Provide Expert Answer\n' +
            'Complete MOLT Replicator guide\n' +
            'Based on official SKILL.md\n' +
            'Calls complete_task()',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 7, 4.55, 7, 4.8)

    # Step 12: Store new learning
    add_box(slide, 5.6, 4.8, 2.9, 0.55,
            'Step 12: Store Learning\n' +
            'query_store.py: store_query()\n' +
            'Saves: query + [molt-replicator]\n' +
            'Next time = instant expert!',
            RGBColor(220, 255, 220), 8.5, False)

    # Explanation boxes
    explain1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(4.5), Inches(0.55))
    explain1.text_frame.text = "🎯 Dynamic Skill Loading: Claude intelligently decides when to fetch skills. Simple queries = no fetch. Complex queries = fetch on-demand. This is the 'fetch_skill' tool in action - available to Claude just like list_databases."
    explain1.text_frame.paragraphs[0].font.size = Pt(8.5)
    explain1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    explain1.text_frame.word_wrap = True

    explain2 = slide.shapes.add_textbox(Inches(5), Inches(5.5), Inches(4.5), Inches(0.55))
    explain2.text_frame.text = "💡 Tool Discovery: fetch_skill is defined in get_agent_tools() alongside MCP tools. Claude sees: {'name': 'fetch_skill', 'description': 'Fetch a CockroachDB best practice skill', 'parameters': {'skill_name': 'string'}}. Claude chooses when to call it!"
    explain2.text_frame.paragraphs[0].font.size = Pt(8.5)
    explain2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    explain2.text_frame.word_wrap = True

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(6.2), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "Skills Table: ai_demo.skills contains 25+ SKILL.md files. Each is a CockroachDB best practice document (2-10 KB each). Fetched only when needed."
    legend.text_frame.paragraphs[0].font.size = Pt(9)
    legend.text_frame.paragraphs[0].font.italic = True
    legend.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

    return slide

def create_presentation():
    """Create presentation with 3 detailed slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create the 3 detailed slides
    create_slide_1_new_query_detailed(prs)
    create_slide_2_learned_query_detailed(prs)
    create_slide_3_complex_skill_fetch(prs)

    # Save
    filename = '/var/www/ai/aidemo2/docs/timbobwith3slides.pptx'
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\nSlide Details:")
    print("  Slide 1: New Query (Cold Start) - Simple query, zero skills")
    print("  Slide 2: Learned Query (Warm Start) - Similarity match found")
    print("  Slide 3: Complex Query - No match, dynamic skill fetch")

    return filename

if __name__ == "__main__":
    create_presentation()
