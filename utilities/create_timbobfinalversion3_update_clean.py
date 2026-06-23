#!/usr/bin/env python3
"""
Final Visual PowerPoint Presentation - Version 3 Update (Clean Build)
Complete presentation with updated slides 6-8
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def add_box(slide, x, y, width, height, text, fill_color, text_size=10, bold=False):
    """Add a colored box with BLACK text"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(2)

    box.text_frame.text = text
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = 1

    # Set ALL paragraphs to black (multi-line text creates multiple paragraphs)
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # BLACK TEXT!

    return box

def add_arrow(slide, x1, y1, x2, y2, label=None):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(2)
    connector.line.end_arrow_type = 2

    if label:
        label_box = slide.shapes.add_textbox(
            Inches((x1 + x2) / 2 - 0.4), Inches((y1 + y2) / 2 - 0.15),
            Inches(0.8), Inches(0.25)
        )
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.italic = True
        label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        label_box.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return connector

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title

    p = title_frame.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_bullet_slide(prs, title, bullets):
    """Add a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    tf = body_shape.text_frame
    tf.clear()

    for bullet in bullets:
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 - level * 2)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(10)

    return slide

def add_architecture_overview(prs):
    """MCP-centric architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text = "Architecture: MCP Protocol + SKILL.md Knowledge Base"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # User
    add_box(slide, 4.25, 0.8, 1.5, 0.5, "User", RGBColor(173, 216, 230))
    add_arrow(slide, 5, 1.3, 5, 1.8)

    # Agent (center)
    add_box(slide, 3.5, 1.8, 3, 0.7, "Agent (agent.py)\nOrchestrator", RGBColor(255, 200, 100), 13)

    # Left - Skill System
    add_box(slide, 0.5, 3, 2.5, 0.7, "SKILL.md System\n25+ Best Practice Guides\nStored in ai_demo.skills",
            RGBColor(255, 240, 200), 10)
    add_arrow(slide, 1.75, 2.5, 4.5, 2.2, "fetch guides")

    # Right - Claude API
    add_box(slide, 7, 3, 2.5, 0.7, "Claude Sonnet 4.6\n(Vertex AI)\nTool Use API",
            RGBColor(200, 230, 255), 11)
    add_arrow(slide, 6.5, 2.2, 8, 3, "API calls")

    # Bottom left - Query Learning (OPTIMIZATION)
    add_box(slide, 0.5, 4.2, 2.5, 0.6, "Query Learning\n(OPTIMIZATION)\n768-dim embeddings",
            RGBColor(220, 255, 220), 10)
    add_arrow(slide, 1.75, 3.7, 1.75, 4.2)

    # Bottom center - MCP Client (MAIN FEATURE)
    add_box(slide, 3.5, 4.2, 2.5, 0.6, "MCP Client\n(MAIN PROTOCOL)\nDynamic tool discovery",
            RGBColor(255, 210, 210), 10, True)
    add_arrow(slide, 5, 2.5, 5, 4.2, "execute")

    # Bottom right - MCP Server
    add_box(slide, 6.5, 4.2, 2.5, 0.6, "MCP Server\ncockroachlabs.cloud/mcp\n11 Database Tools",
            RGBColor(240, 220, 255), 10)
    add_arrow(slide, 6.5, 2.5, 7.75, 4.2, "discover")

    # CockroachDB Cluster (THE TARGET)
    add_box(slide, 2.5, 5.5, 5, 0.9, "CockroachDB Cluster\n(THE CLUSTER BEING MANAGED)\nAlso stores: skills, query history, embeddings",
            RGBColor(100, 150, 200), 11, True)
    add_arrow(slide, 5, 4.8, 5, 5.5)

    # Key message
    msg = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.6))
    msg.text_frame.text = "Core Value: MCP protocol provides 11 database tools + SKILL.md files provide expert guidance\nOptimization: Query learning pre-loads relevant skills (saves 30-40% tokens)"
    for p in msg.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_slide_6_no_skill_needed(prs):
    """Slide 6: Simple task - Claude doesn't need SKILL.md"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 1: Simple Task - No SKILL.md Needed"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1: User
    box1 = add_box(slide, 0.5, 1, 2, 0.6, 'User Query\n"List all databases"',
            RGBColor(173, 216, 230), 11, True)

    # Step 2: MCP Discovery
    box2 = add_box(slide, 3, 1, 2, 0.6,
            'MCP Discovery\nmcp_client.py\nlist_tools()\n11 tools found',
            RGBColor(255, 210, 210), 9.5, False)
    add_arrow(slide, 2.5, 1.3, 3, 1.3)

    # Step 3: Check History
    box3 = add_box(slide, 5.5, 1, 2, 0.6,
            'Check History\nquery_store.py\nfind_similar_queries()\nNo match',
            RGBColor(255, 230, 230), 9.5, False)
    add_arrow(slide, 5, 1.3, 5.5, 1.3)

    # Step 4: Call Claude (no skills)
    box4 = add_box(slide, 7.75, 1, 2, 0.6,
            'Call Claude API\nagent.py\nNo skills pre-loaded\nJust tools available',
            RGBColor(200, 230, 255), 9.5, False)
    add_arrow(slide, 7.5, 1.3, 7.75, 1.3)

    # Step 5: Claude analyzes
    box5 = add_box(slide, 7.75, 2, 2, 0.7,
            'Claude Decides\n"This is simple"\n"I know list_databases"\n"No skill needed"',
            RGBColor(220, 255, 220), 9.5, True)
    add_arrow(slide, 8.75, 1.6, 8.75, 2)

    # Step 6: Execute MCP
    box6 = add_box(slide, 5.5, 2, 2, 0.7,
            'Execute MCP Tool\nmcp_client.py\nexecute_tool_call(\n"list_databases", {})',
            RGBColor(255, 210, 210), 9.5, False)
    add_arrow(slide, 7.75, 2.35, 7.5, 2.35)

    # Step 7: Results
    box7 = add_box(slide, 3, 2, 2, 0.7,
            'Database Results\ndefaultdb\nai_demo\npostgres\nsystem',
            RGBColor(180, 220, 180), 9.5, False)
    add_arrow(slide, 5.5, 2.35, 5, 2.35)

    # Step 8: Claude formats
    box8 = add_box(slide, 0.5, 2, 2, 0.7,
            'Claude Formats\n"Your cluster has 4\ndatabases:"\n(lists them)',
            RGBColor(255, 240, 200), 9.5, False)
    add_arrow(slide, 3, 2.35, 2.5, 2.35)

    # Step 9: Store Learning
    box9 = add_box(slide, 0.5, 3.1, 2, 0.65,
            'Store Query\nquery_store.py\nstore_query()\nskills_used: []',
            RGBColor(220, 255, 220), 9.5, False)
    add_arrow(slide, 1.5, 2.7, 1.5, 3.1)

    # Key insight box
    insight = slide.shapes.add_textbox(Inches(3), Inches(3.1), Inches(6.75), Inches(0.65))
    insight.text_frame.text = "Key Point: Claude uses its training + tool descriptions to handle simple tasks.\nNo SKILL.md needed! Stored with empty skills_used[] for future reference."
    for p in insight.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Flow summary
    summary = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.3))
    summary.text_frame.text = """Complete Flow Summary:

1. User asks simple question
2. MCP discovers 11 tools
3. Check history (no match)
4. Call Claude with tools only (no skills)
5. Claude: "I know this, use list_databases"
6. Execute MCP tool
7. Get results from CockroachDB
8. Claude formats answer
9. Store: query + embedding + skills_used:[] (empty - no skills needed!)

Stored in ai_demo.query_history:
   - query_text: "List all databases"
   - embedding: [768-dimensional vector]
   - skills_used: []  (Empty because Claude didn't need help!)
   - timestamp: 2026-04-05 10:23:15

Next time similar query: Still no skills loaded (because none were used)
Simple queries stay fast and lightweight!"""
    for p in summary.text_frame.paragraphs:
        p.font.size = Pt(9.5)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_slide_7_claude_fetches_skill(prs):
    """Slide 7: Complex task - Claude calls fetch_skill, stores result"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 2: Complex Task - Claude Fetches SKILL.md (No History Match)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    box1 = add_box(slide, 0.4, 0.7, 2, 0.6, 'User Query\n"Set up MOLT\nReplicator"',
            RGBColor(173, 216, 230), 11, True)

    # Step 2
    box2 = add_box(slide, 2.7, 0.7, 1.8, 0.6,
            'MCP Discovery\n11 tools\nincluding\nfetch_skill',
            RGBColor(255, 210, 210), 9, False)
    add_arrow(slide, 2.4, 1, 2.7, 1)

    # Step 3
    box3 = add_box(slide, 4.8, 0.7, 1.8, 0.6,
            'Check History\nfind_similar()\nNo match\n(never asked)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 4.5, 1, 4.8, 1)

    # Step 4
    box4 = add_box(slide, 7, 0.7, 2.6, 0.6,
            'Call Claude\nNo skills pre-loaded\nTools: MCP + fetch_skill\nUser query included',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 6.6, 1, 7, 1)

    # Step 5
    box5 = add_box(slide, 7, 1.7, 2.6, 0.7,
            'Claude Analyzes\n"MOLT migration is complex"\n"I need official guide"\nDecides: call fetch_skill',
            RGBColor(255, 240, 150), 9, True)
    add_arrow(slide, 8.3, 1.3, 8.3, 1.7)

    # Step 6
    box6 = add_box(slide, 4.8, 1.7, 1.8, 0.7,
            'Claude Calls\nfetch_skill(\n"onboarding-\nand-migrations/\nmolt-replicator")',
            RGBColor(255, 200, 150), 8.5, True)
    add_arrow(slide, 7, 2.05, 6.6, 2.05)

    # Step 7
    box7 = add_box(slide, 2.7, 1.7, 1.8, 0.7,
            'Fetch SKILL.md\nskill_fetcher.py\nQuery ai_demo\n.skills table\nReturns 8KB',
            RGBColor(255, 240, 200), 8.5, False)
    add_arrow(slide, 4.8, 2.05, 4.5, 2.05)

    # Step 8
    box8 = add_box(slide, 0.4, 1.7, 2, 0.7,
            'SKILL.md Loaded\nBest practices\nSetup steps\nCommon issues',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 2.7, 2.05, 2.4, 2.05)

    # Step 9
    box9 = add_box(slide, 0.4, 2.8, 2, 0.65,
            'Re-process\nClaude now has:\nTools\nSKILL.md\nUser query',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 1.4, 2.4, 1.4, 2.8)

    # Step 10
    box10 = add_box(slide, 2.7, 2.8, 1.8, 0.65,
            'Expert Answer\nClaude uses\nSKILL.md to\nguide setup',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.4, 3.1, 2.7, 3.1)

    # Step 11
    box11 = add_box(slide, 4.8, 2.8, 2.2, 0.65,
            'Store Learning\nquery + embedding +\nskills_used:\n["molt-replicator"]\nquery_history',
            RGBColor(180, 255, 180), 8.5, True)
    add_arrow(slide, 4.5, 3.1, 4.8, 3.1)

    # Key insight
    insight = slide.shapes.add_textbox(Inches(0.4), Inches(3.7), Inches(9.2), Inches(0.6))
    insight.text_frame.text = 'Critical Learning Moment: Query stored with skills_used:["onboarding-and-migrations/molt-replicator"]\nNext similar query will pre-load this skill automatically!'
    for p in insight.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 0, 102)

    # Detailed storage
    storage = slide.shapes.add_textbox(Inches(0.4), Inches(4.5), Inches(9.2), Inches(1.9))
    storage.text_frame.text = """What Gets Stored in ai_demo.query_history:

Database Row Created:
   - query_text: "Set up MOLT Replicator"
   - embedding: [768-dim vector from sentence-transformers]
   - skills_used: ["onboarding-and-migrations/molt-replicator"]  (THIS IS KEY!)
   - timestamp: 2026-04-05 11:45:22

Why This Matters:
   - Future queries like "help with MOLT setup" or "MOLT migration guide"
   - pgvector similarity search finds this stored query (cosine similarity > 0.3)
   - Agent sees: "Someone asked about MOLT before and used molt-replicator skill"
   - Pre-loads that skill BEFORE calling Claude - Saves 1 API call!

Claude made the decision: "I need molt-replicator skill"
Agent just executed and remembered it for next time."""
    for p in storage.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_slide_8_preloaded_from_history(prs):
    """Slide 8: Similar query - skill pre-loaded from history"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 3: Optimization - Skill Pre-loaded from Query History"
    title.text_frame.paragraphs[0].font.size = Pt(26)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    box1 = add_box(slide, 0.4, 0.7, 2, 0.6, 'User Query\n"Help with MOLT\nmigration"',
            RGBColor(173, 216, 230), 11, True)

    # Step 2
    box2 = add_box(slide, 2.7, 0.7, 1.8, 0.6,
            'Generate\nEmbedding\n768-dim\nvector',
            RGBColor(255, 230, 230), 9, False)
    add_arrow(slide, 2.4, 1, 2.7, 1)

    # Step 3 - MATCH FOUND!
    box3 = add_box(slide, 4.8, 0.7, 2.2, 0.6,
            'Check History\npgvector search\nMATCH FOUND!\nSimilarity: 0.82',
            RGBColor(180, 255, 180), 9, True)
    add_arrow(slide, 4.5, 1, 4.8, 1)

    # Step 4
    box4 = add_box(slide, 7.3, 0.7, 2.3, 0.6,
            'Found Match:\n"Set up MOLT Replicator"\nskills_used:\n["molt-replicator"]',
            RGBColor(220, 255, 220), 8.5, False)
    add_arrow(slide, 7, 1, 7.3, 1)

    # Step 5
    box5 = add_box(slide, 7.3, 1.7, 2.3, 0.7,
            'Pre-load Skill\nskill_fetcher.py\nget_skill(\n"molt-replicator")\nBEFORE Claude call',
            RGBColor(255, 240, 200), 8.5, True)
    add_arrow(slide, 8.45, 1.3, 8.45, 1.7)

    # Step 6
    box6 = add_box(slide, 4.8, 1.7, 2.2, 0.7,
            'Call Claude\nWith skill\nPRE-LOADED\n+ tools\n+ query',
            RGBColor(200, 230, 255), 9, True)
    add_arrow(slide, 7.3, 2.05, 7, 2.05)

    # Step 7
    box7 = add_box(slide, 2.7, 1.7, 1.8, 0.7,
            'Claude Has\nEverything!\nNo need to\ncall fetch_skill',
            RGBColor(220, 255, 220), 9, True)
    add_arrow(slide, 4.8, 2.05, 4.5, 2.05)

    # Step 8
    box8 = add_box(slide, 0.4, 1.7, 2, 0.7,
            'Expert Answer\nImmediately\nUsing SKILL.md\nguidance',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.7, 2.05, 2.4, 2.05)

    # Step 9
    box9 = add_box(slide, 0.4, 2.8, 2, 0.65,
            'Store Query\nConfirm skill\nhelped again\nReinforces',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 1.4, 2.4, 1.4, 2.8)

    # Performance comparison
    perf = slide.shapes.add_textbox(Inches(2.7), Inches(2.8), Inches(6.9), Inches(0.65))
    perf.text_frame.text = "Performance: 1 Claude API call (vs 2 in Flow 2)\nCost Savings: ~40% fewer tokens | Time Savings: ~1 second faster"
    for p in perf.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Detailed comparison
    comparison = slide.shapes.add_textbox(Inches(0.4), Inches(3.7), Inches(9.2), Inches(2.7))
    comparison.text_frame.text = """What Happened Behind the Scenes:

Step 3 Details - pgvector Similarity Search:
   - Query "Help with MOLT migration" generates embedding
   - CockroachDB executes vector similarity search with HNSW index
   - HNSW index makes this sub-millisecond!
   - Best match: "Set up MOLT Replicator" (similarity 0.82)
   - Extract skills_used: ["onboarding-and-migrations/molt-replicator"]

Step 5 - Optimization in Action:
   - Agent: "I'll load molt-replicator skill BEFORE calling Claude"
   - Fetches 8KB SKILL.md from ai_demo.skills table
   - Includes it in Claude API call's system prompt

The Result:
   Flow 2 (no history): User to Claude to fetch_skill to Claude again to Answer (2 API calls)
   Flow 3 (with history): User to [pre-load skill] to Claude to Answer (1 API call)

This is the Learning Optimization: Same quality answer, half the API calls!"""
    for p in comparison.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_real_world_use_cases(prs):
    """MCP-focused use cases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Real-World Use Cases: Managing CockroachDB Cloud"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Use cases
    add_box(slide, 0.5, 1.2, 4.2, 0.6, "Cluster Operations", RGBColor(200, 230, 255), 16, True)
    uc1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(1.2))
    uc1.text_frame.text = """Show me cluster info
List all databases
Check cluster health
Uses: get_cluster, list_databases
SKILL.md: operations guides"""
    for p in uc1.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.3, 1.2, 4.2, 0.6, "Database Creation", RGBColor(220, 255, 220), 16, True)
    uc2 = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(1.2))
    uc2.text_frame.text = """Create test_db for staging
Create users table with schema
Uses: create_database, create_table
With safety confirmations
SKILL.md: schema design best practices"""
    for p in uc2.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.5, 3.4, 4.2, 0.6, "Query and Analysis", RGBColor(255, 230, 230), 16, True)
    uc3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(4.2), Inches(1.2))
    uc3.text_frame.text = """Show slow queries
Analyze range distribution
Find table hotspots
Uses: select_query with expert SQL
SKILL.md: observability guides"""
    for p in uc3.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.3, 3.4, 4.2, 0.6, "Migrations", RGBColor(255, 240, 200), 16, True)
    uc4 = slide.shapes.add_textbox(Inches(5.3), Inches(4.1), Inches(4.2), Inches(1.2))
    uc4.text_frame.text = """Set up MOLT Replicator
Migrate from MySQL
Best practices for MOLT Fetch
Uses: MCP tools + extensive SKILL.md
Guides through complex processes"""
    for p in uc4.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 2.9, 5.6, 4.2, 0.6, "Learning and Support", RGBColor(240, 220, 255), 16, True)
    uc5 = slide.shapes.add_textbox(Inches(2.9), Inches(6.3), Inches(4.2), Inches(0.9))
    uc5.text_frame.text = """How do I optimize this query?
Security best practices?
Pulls from 25+ SKILL.md files
Consistent expert answers 24/7"""
    for p in uc5.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_why_cockroachdb(prs):
    """Why CockroachDB - it's the cluster being managed!"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Why CockroachDB: The Cluster Being Managed"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Main message
    main = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.7))
    main.text_frame.text = "This demo manages a CockroachDB Cloud cluster using MCP protocol.\nThe same database ALSO stores the agent's metadata (skills, query history)."
    for p in main.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 0, 102)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Features
    add_box(slide, 0.5, 2, 4.3, 0.7, "Target Cluster", RGBColor(255, 210, 210), 14, True)
    f1 = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.3), Inches(1.3))
    f1.text_frame.text = """The cluster being managed
MCP provides tools to interact
Create databases, tables
Run queries and analysis
Monitor health and performance"""
    for p in f1.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.2, 2, 4.3, 0.7, "Metadata Storage", RGBColor(255, 240, 200), 14, True)
    f2 = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4.3), Inches(1.3))
    f2.text_frame.text = """ai_demo.skills: 25+ SKILL.md files
ai_demo.query_history: Learning data
Embeddings with pgvector
All in the same cluster!
Eating its own dog food"""
    for p in f2.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.5, 4.4, 4.3, 0.7, "Production Features", RGBColor(220, 255, 220), 14, True)
    f3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(4.3), Inches(1.3))
    f3.text_frame.text = """ACID transactions for writes
Strong consistency guarantees
OAuth authentication
Distributed globally
Horizontal scalability"""
    for p in f3.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.2, 4.4, 4.3, 0.7, "Vector Capabilities", RGBColor(200, 230, 255), 14, True)
    f4 = slide.shapes.add_textbox(Inches(5.2), Inches(5.2), Inches(4.3), Inches(1.3))
    f4.text_frame.text = """Native pgvector support
768-dim embeddings
HNSW index for fast search
No external vector DB needed
SQL + vectors in one place"""
    for p in f4.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Bottom message
    bottom = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.6))
    bottom.text_frame.text = "Unique Value: CockroachDB is both the TARGET (cluster being managed) and the INFRASTRUCTURE (stores skills, learning data). One distributed database for everything!"
    for p in bottom.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return slide

def create_presentation():
    """Create complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs, "CockroachDB Cloud Management Agent",
                   "MCP Protocol + SKILL.md Knowledge Base")

    # 2. Overview - MCP FIRST!
    overview = [
        "What This Demo Does",
        ("Manages CockroachDB Cloud cluster via MCP protocol", 1),
        ("Create databases/tables, run queries, monitor health, all via natural language", 2),
        ("11 database tools discovered dynamically", 1),
        ("Agent learns available operations from MCP server at startup - not hardcoded", 2),
        ("25+ SKILL.md files provide expert guidance", 1),
        ("Best practices for migrations, performance tuning, security, troubleshooting", 2),
        ("Query learning optimizes skill selection", 1),
        ("Pre-loads relevant skills from past similar queries - saves time and tokens", 2),
        "",
        "Core Components",
        ("MCP Protocol: Connects to CockroachDB Cloud", 1),
        ("Standard protocol eliminates custom integration - secure, reliable communication", 2),
        ("Dynamic Tool Discovery: list_databases, create_table, etc.", 1),
        ("Agent adapts automatically when new database features are added - future-proof", 2),
        ("SKILL.md System: Best practices from GitHub docs", 1),
        ("Official CockroachDB expertise built in - always up-to-date with latest docs", 2),
        ("fetch_skill: Claude calls when expertise needed", 1),
        ("Claude intelligently requests guidance only for complex tasks - efficient context use", 2),
        ("Query Learning: OPTIMIZATION - pre-loads relevant skills", 1),
        ("Remembers which skills helped before - skips fetch_skill step for faster responses", 2),
    ]
    add_bullet_slide(prs, "Overview", overview)

    # 3. How to Use the Demo
    setup = [
        "Prerequisites",
        ("CockroachDB Cloud cluster (free tier works!)", 1),
        ("Cluster ID and API key from CockroachDB Cloud", 1),
        ("Google Cloud project with Vertex AI API enabled", 1),
        ("Python 3.8+ with pip", 1),
        "",
        "Setup Steps",
        ("1. Copy config.template.json to config.json", 1),
        ("2. Add your CockroachDB cluster_id and api_key", 1),
        ("3. Add your Google Cloud project_id and region", 1),
        ("4. Install dependencies: pip install -r requirements.txt", 1),
        ("5. Run setup: python setup_database.py", 1),
        ("   - Creates ai_demo database", 1),
        ("   - Creates query_history and skills tables", 1),
        ("   - Loads 25+ SKILL.md files from GitHub", 1),
        "",
        "Running the Demo",
        ("python main.py (interactive mode)", 1),
        ('Ask questions: "list databases", "check cluster health"', 1),
        ("Type 'quit' to exit", 1),
    ]
    add_bullet_slide(prs, "How to Use This Demo", setup)

    # 4. Section: Architecture
    add_section_slide(prs, "Architecture")

    # 5. Architecture diagram
    add_architecture_overview(prs)

    # 6-8. UPDATED FLOW DIAGRAMS
    add_slide_6_no_skill_needed(prs)
    add_slide_7_claude_fetches_skill(prs)
    add_slide_8_preloaded_from_history(prs)

    # 9. Section: Use Cases
    add_section_slide(prs, "Real-World Applications")

    # 10. Use cases
    add_real_world_use_cases(prs)

    # 11. Why CockroachDB
    add_why_cockroachdb(prs)

    # 12. Section: Technical
    add_section_slide(prs, "Technical Details")

    # 13. MCP Protocol details
    mcp_details = [
        "MCP Protocol (Model Context Protocol)",
        ("Standard protocol for AI-to-database communication", 1),
        ("JSON-RPC 2.0 over HTTPS", 1),
        ("Server: cockroachlabs.cloud/mcp", 1),
        ("Tools/list: Discovers available operations", 1),
        ("Tools/call: Executes database operations", 1),
        "",
        "11 Discovered MCP Tools",
        ("Read: list_databases, list_tables, select_query", 1),
        ("Write: create_database, create_table, insert_rows", 1),
        ("Cluster: get_cluster, get_cluster_cert, get_cluster_api_key", 1),
        ("Plus 2 more... all discovered at runtime!", 1),
        "",
        "SKILL.md Knowledge Base",
        ("25+ expert guides from cockroachlabs/cockroach GitHub", 1),
        ("Categories: operations, security, migrations, diagnostics", 1),
        ("Stored in CockroachDB ai_demo.skills table", 1),
        ("fetch_skill tool: Claude loads when needed", 1),
    ]
    add_bullet_slide(prs, "MCP Protocol & SKILL.md System", mcp_details)

    # 14. Learning optimization
    learning_details = [
        "Query Learning: Optional Optimization",
        ("NOT the main feature - just makes it faster!", 1),
        ("Works perfectly without learning (just 1 extra API call)", 1),
        "",
        "How It Works:",
        ("Every query to 768-dim embedding (sentence-transformers)", 1),
        ("Store: query text + embedding + skills_used[]", 1),
        ("Next time: pgvector finds similar (HNSW index)", 1),
        ("Pre-load those skills to Skip fetch_skill to Faster", 1),
        "",
        "Performance Impact:",
        ("Without learning: ~3 sec (2 Claude API calls)", 1),
        ("With learning (1st): ~3 sec + store", 1),
        ("With learning (2nd+): ~2 sec (1 Claude API call)", 1),
        ("Saves: 30-40% tokens, 1 API call, ~1 second", 1),
        "",
        "Why It Matters:",
        ("Means to efficiently use SKILL.md files", 1),
        ("Not the goal - just an optimization", 1),
        ("Demo value: MCP + SKILL.md, not learning", 1),
    ]
    add_bullet_slide(prs, "Query Learning (Optimization)", learning_details)

    # 15. Write operations
    write_ops = [
        "Write Operations (April 5, 2026)",
        ("create_database: Create new databases", 1),
        ("create_table: Create tables with DDL", 1),
        ("insert_rows: Insert data safely", 1),
        "",
        "Safety Features",
        ("User confirmation REQUIRED for every write", 1),
        ("Shows exact operation before execution", 1),
        ("Can cancel any operation", 1),
        ("OAuth authentication", 1),
        ("Dark orange warning for visibility", 1),
        "",
        "Not Supported (By Design)",
        ("UPDATE, DELETE - too risky", 1),
        ("DROP, TRUNCATE - destructive", 1),
        ("All dangerous operations blocked", 1),
    ]
    add_bullet_slide(prs, "Write Operations with Safety", write_ops)

    # 16. Summary
    summary = [
        "Core Demo Value",
        ("MCP Protocol: Standard interface to CockroachDB Cloud", 1),
        ("SKILL.md Files: 25+ expert guides from official docs", 1),
        ("Natural language to Database operations", 1),
        ("Complex tasks guided by best practices", 1),
        "",
        "Optimizations",
        ("Dynamic tool discovery: Not brittle, adapts automatically", 1),
        ("Query learning: Pre-loads relevant skills (30-40% faster)", 1),
        ("Zero-skill mode: Simple queries need no SKILL.md", 1),
        "",
        "Why CockroachDB",
        ("It's the cluster being managed (via MCP)", 1),
        ("Also stores: SKILL.md files, query history, embeddings", 1),
        ("One database for everything: target + infrastructure", 1),
        ("pgvector: No external vector DB needed", 1),
        ("ACID transactions: Safe writes, no lost learning", 1),
        ("Global distribution: Deploy agents worldwide", 1),
    ]
    add_bullet_slide(prs, "Summary: MCP + SKILL.md = Powerful Agent", summary)

    # Save
    filename = '/var/www/ai/aidemo2/docs/timbobfinalversion3-update.pptx'
    prs.save(filename)
    print(f"Presentation created: {filename}")
    print(f"Total slides: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    create_presentation()
