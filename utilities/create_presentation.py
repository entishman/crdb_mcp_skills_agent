#!/usr/bin/env python3
"""
Create PowerPoint presentation for AIdemo2 project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_points):
    """Add a bullet point slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title

    tf = body_shape.text_frame
    tf.clear()

    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_two_column_slide(prs, title, left_points, right_points):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title

    # Remove default content placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape:
            sp = shape.element
            sp.getparent().remove(sp)

    # Add left text box
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    tf = left.text_frame
    for point in left_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)

    # Add right text box
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    tf = right.text_frame
    for point in right_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)

    return slide

# Slide 1: Title
add_title_slide(
    prs,
    "AIdemo2: Intelligent CockroachDB Agent",
    "Self-Learning AI Assistant with Query-Based Learning"
)

# Slide 2: The Problem
add_content_slide(
    prs,
    "The Problem",
    [
        "Traditional database assistants are static",
        "• Load ALL documentation every time (wasteful)",
        "• Don't learn from interactions",
        "• Can't understand paraphrased questions",
        "",
        "Result: Slow, expensive, not getting smarter"
    ]
)

# Slide 3: The Solution
add_content_slide(
    prs,
    "Our Solution: Query-Based Learning",
    [
        "✓ Learns from every interaction",
        "✓ Finds similar past queries using AI embeddings",
        "✓ Loads only relevant documentation (2-5 files vs 25)",
        "✓ Understands meaning, not just keywords",
        "",
        "Result: 95% token savings, gets smarter over time"
    ]
)

# Slide 4: How It Works
add_content_slide(
    prs,
    "How It Works: The Learning Loop",
    [
        "1. User asks: 'Are there hotspots in testdb?'",
        "2. Generate semantic embedding (384-dim vector)",
        "3. Search for similar past queries (cosine similarity)",
        "4. Load skills that helped with similar queries",
        "5. Claude answers using MCP tools",
        "6. Store: query + skills that actually helped",
        "",
        "Next similar query → Pre-load learned skills!"
    ]
)

# Slide 5: Architecture
add_content_slide(
    prs,
    "System Architecture",
    [
        "User Query",
        "  ↓",
        "sentence-transformers (embeddings)",
        "  ↓",
        "CockroachDB pgvector (similarity search)",
        "  ↓",
        "Load learned skills from database",
        "  ↓",
        "Claude Sonnet 4.6 (via Vertex AI)",
        "  ↓",
        "MCP Server → CockroachDB Cluster",
        "  ↓",
        "Store learning for future queries"
    ]
)

# Slide 6: Key Features
add_two_column_slide(
    prs,
    "Key Features",
    [
        "Smart Features:",
        "• Dynamic skill loading",
        "• Semantic similarity search",
        "• Progressive learning",
        "• Skill usage tracking",
        "",
        "Safety Features:",
        "• Read-only enforcement",
        "• Cannot modify data",
        "• Refuses write operations"
    ],
    [
        "Performance:",
        "• 95% token savings",
        "• ~2 second response time",
        "• Gets faster over time",
        "",
        "Storage:",
        "• ~2KB per query",
        "• 25 skills in database",
        "• HNSW index for speed"
    ]
)

# Slide 7: Technical Stack
add_content_slide(
    prs,
    "Technology Stack",
    [
        "AI & ML:",
        "• Claude Sonnet 4.6 (Anthropic via Vertex AI)",
        "• sentence-transformers (all-MiniLM-L6-v2)",
        "",
        "Database:",
        "• CockroachDB with pgvector extension",
        "• HNSW index for fast similarity search",
        "",
        "Other:",
        "• MCP (Model Context Protocol)",
        "• Python 3.12 (2,545 lines of code)"
    ]
)

# Slide 8: Learning Example
add_content_slide(
    prs,
    "Learning in Action",
    [
        "Query 1: 'How do I check cluster health?'",
        "  → No match found",
        "  → Load 3 default skills",
        "  → Store: ['reviewing-cluster-health']",
        "",
        "Query 2: 'Is my cluster healthy?'",
        "  → Match found! (0.70 similarity)",
        "  → Pre-load: 'reviewing-cluster-health'",
        "  → Answer immediately ✓",
        "",
        "System learned the pattern!"
    ]
)

# Slide 9: Embeddings Upgrade
add_two_column_slide(
    prs,
    "The Embeddings Upgrade",
    [
        "Before (Hash-based):",
        "• 'hotspots in testdb'",
        "  vs",
        "  'hot ranges in testdb'",
        "• Similarity: 0.32",
        "• Missed match ✗",
        "",
        "Problem:",
        "• Different words",
        "• Same meaning",
        "• System can't tell"
    ],
    [
        "After (sentence-transformers):",
        "• 'hotspots in testdb'",
        "  vs",
        "  'hot ranges in testdb'",
        "• Similarity: 0.53",
        "• Found match ✓",
        "",
        "Solution:",
        "• Semantic understanding",
        "• Captures meaning",
        "• Finds paraphrases"
    ]
)

# Slide 10: Performance Metrics
add_content_slide(
    prs,
    "Performance Metrics",
    [
        "Token Efficiency:",
        "• Legacy mode: ~50,000 tokens/query",
        "• Query learning: ~2,500 tokens/query",
        "• Savings: 95% ✓",
        "",
        "Response Time:",
        "• Embedding generation: 50ms",
        "• Similarity search: 5ms",
        "• Claude API: 2,000ms",
        "• Total: ~2.1 seconds",
        "",
        "Storage: ~2KB per query (tiny!)"
    ]
)

# Slide 11: Code Statistics
add_content_slide(
    prs,
    "Codebase",
    [
        "Total: 2,545 lines of Python",
        "",
        "Core Components:",
        "• agent.py: 783 lines (orchestrator)",
        "• query_store.py: 303 lines (learning)",
        "• mcp_tools.py: 309 lines (tool definitions)",
        "• skill_fetcher.py: 268 lines (skill loader)",
        "• mcp_client.py: 247 lines (MCP client)",
        "• ollama_embedding_manager.py: 184 lines (embeddings)",
        "",
        "All heavily documented!"
    ]
)

# Slide 12: Database Schema
add_content_slide(
    prs,
    "Database Schema",
    [
        "ai_demo.query_history:",
        "• query_text: User's question",
        "• query_embedding: VECTOR(384)",
        "• skills_used: TEXT[] - which skills helped",
        "• execution_time_ms: Performance tracking",
        "",
        "ai_demo.skills:",
        "• skill_name: e.g., 'reviewing-cluster-health'",
        "• content: Full SKILL.md markdown",
        "",
        "HNSW index on embeddings for fast search"
    ]
)

# Slide 13: Future Improvements
add_content_slide(
    prs,
    "Future Enhancements",
    [
        "Performance:",
        "• Connection pooling (50% faster)",
        "• Higher quality embeddings (nomic-embed-text)",
        "",
        "Features:",
        "• User feedback loop (rate answers)",
        "• Metrics dashboard",
        "• Multi-user support",
        "",
        "All documented in ARCHITECTURE.md"
    ]
)

# Slide 14: Summary
add_content_slide(
    prs,
    "Summary",
    [
        "Built a self-learning AI assistant that:",
        "",
        "✓ Learns from every interaction",
        "✓ Uses semantic embeddings (understands meaning)",
        "✓ Loads only needed documentation (95% savings)",
        "✓ Gets smarter over time",
        "✓ Safe by design (read-only)",
        "",
        "2,545 lines of intelligent, documented code",
        "Ready for production!"
    ]
)

# Save presentation
output_file = "AIdemo2_Presentation.pptx"
prs.save(output_file)
print(f"✓ Presentation created: {output_file}")
print(f"  Total slides: {len(prs.slides)}")
