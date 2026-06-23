#!/usr/bin/env python3
"""
Create Architecture Slide - Single PowerPoint slide
Based on arch1.png but updated for current MCP-first architecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def add_box(slide, x, y, width, height, text, fill_color, text_size=10, bold=False):
    """Add a colored box with BLACK text"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(2)

    box.text_frame.text = text
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = 1

    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0, 0, 0)

    return box

def add_arrow(slide, x1, y1, x2, y2, color=RGBColor(0, 0, 0), width=2):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    connector.line.end_arrow_type = 2
    return connector

def create_architecture_slide():
    """Create single architecture slide"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5))
    title.text = "Aidemo2 Complete Architecture"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.3), Inches(0.5), Inches(9.4), Inches(0.3))
    subtitle.text = "Self-Learning CockroachDB Agent with Query-Based Learning"
    subtitle.text_frame.paragraphs[0].font.size = Pt(14)
    subtitle.text_frame.paragraphs[0].font.italic = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 60, 60)
    subtitle.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # USER at top
    add_box(slide, 4.25, 0.9, 1.5, 0.4, "User\nNatural Language Query",
            RGBColor(173, 216, 230), 10, True)
    add_arrow(slide, 5, 1.3, 5, 1.6)

    # AGENT - Central orchestrator
    add_box(slide, 3.25, 1.6, 3.5, 0.6, "Agent (agent.py)\nOrchestrator & Learning Engine",
            RGBColor(255, 200, 100), 11, True)

    # Left side - LEARNING COMPONENTS (Optimization)
    add_box(slide, 0.3, 2.5, 1.8, 0.55, "Query Store\n(query_store.py)\nSimilarity Search",
            RGBColor(255, 230, 230), 9, False)
    add_arrow(slide, 2.1, 2.75, 3.25, 1.95, RGBColor(128, 0, 128), 2)

    add_box(slide, 0.3, 3.2, 1.8, 0.55, "Embedding Manager\n(ollama_embedding_manager.py)\n768-dim Vectors",
            RGBColor(255, 230, 230), 9, False)
    add_arrow(slide, 2.1, 3.45, 3.25, 2.1, RGBColor(128, 0, 128), 2)

    # Ollama LLM - external service for embeddings
    add_box(slide, 0.3, 3.9, 1.8, 0.45, "Ollama LLM\n(Local/Remote)\nall-mpnet-base-v2\nEmbedding Generation",
            RGBColor(220, 220, 255), 7.5, False)
    add_arrow(slide, 1.2, 3.75, 1.2, 3.9, RGBColor(128, 0, 128), 1)

    # Center bottom - COCKROACHDB MANAGER
    add_box(slide, 3.25, 2.5, 1.75, 0.55, "CockroachDB Manager\n(database interaction)\nread/write queries",
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 5, 2.2, 5, 2.5, RGBColor(0, 0, 255), 2)

    # Right side - MCP PROTOCOL (Main Feature)
    add_box(slide, 5.2, 2.5, 1.6, 0.55, "MCP Client\n(mcp_client.py)\nJSON-RPC 2.0",
            RGBColor(255, 210, 210), 9, True)
    add_arrow(slide, 5.2, 2.75, 6.75, 1.95, RGBColor(255, 0, 0), 3)

    # CockroachDB MCP Server
    add_box(slide, 5.2, 3.2, 1.6, 0.55, "CockroachDB\nMCP Server\ncockroachlabs.cloud/mcp\nJSON-RPC + SSE",
            RGBColor(240, 200, 200), 7.5, True)
    add_arrow(slide, 6, 3.05, 6, 3.2, RGBColor(255, 0, 0), 2)

    add_box(slide, 7, 2.5, 1.6, 0.55, "SKILL.md Files\n(skill_fetcher.py)\n25+ guides",
            RGBColor(255, 240, 200), 9, True)
    add_arrow(slide, 7, 2.75, 6.75, 2.1, RGBColor(255, 140, 0), 2)

    # CLAUDE VERTEX AI
    add_box(slide, 8.8, 1.5, 1, 0.7, "Claude\n(Vertex AI)\nSonnet 4.6\nTool Use API",
            RGBColor(200, 230, 255), 9, True)
    add_arrow(slide, 6.75, 1.9, 8.8, 1.85, RGBColor(0, 100, 0), 2)

    # External services - GitHub
    add_box(slide, 8.8, 3.2, 1, 0.45, "GitHub\ncockroachdb/\ncockroach\n25 SKILL.md",
            RGBColor(240, 240, 240), 7.5, False)
    add_arrow(slide, 8.8, 3.2, 7.8, 3.05, RGBColor(100, 100, 100), 1)

    # COCKROACHDB CLUSTER at bottom
    cluster_box = add_box(slide, 0.3, 4.2, 9.4, 0.5, "CockroachDB Cluster (ttc-demo)",
            RGBColor(100, 150, 200), 12, True)

    # Arrows from components to cluster
    add_arrow(slide, 1.2, 4.35, 1.5, 4.2, RGBColor(128, 0, 128), 2)
    add_arrow(slide, 4.1, 3.05, 4.1, 4.2, RGBColor(0, 0, 255), 2)
    add_arrow(slide, 6, 3.75, 5.5, 4.2, RGBColor(255, 0, 0), 2)

    # Database boxes inside cluster
    add_box(slide, 0.5, 4.9, 2.8, 0.9, "ai_demo database\n(Learning Storage)\n\nquery_history table:\n- query_text\n- embedding (768-dim)\n- skills_used[]\n\nskills table:\n- skill_name\n- content (SKILL.md)",
            RGBColor(180, 220, 180), 7.5, False)

    add_box(slide, 3.5, 4.9, 2.8, 0.9, "testdb\n(User Data)\n\nApplication tables\nUser-created schemas\nTest data",
            RGBColor(180, 220, 180), 8, False)

    add_box(slide, 6.5, 4.9, 2.8, 0.9, "defaultdb\n(System DB)\n\nCockroachDB metadata\nSystem tables",
            RGBColor(180, 220, 180), 8, False)

    # Legend at bottom
    legend_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(9.4), Inches(1.3))
    legend_text = """Legend:
User Interface: Light blue  |  Agent: Orange (center)  |  Learning (optimization): Pink  |  MCP Protocol: Red (main feature)
SKILL.md: Yellow (main feature)  |  Claude/LLM: Light blue  |  External Services: Gray  |  CockroachDB: Green

Key Components:
• Claude (Vertex AI): Sonnet 4.6 with tool use API - Makes intelligent decisions about which tools/skills to use
• CockroachDB MCP Server: cockroachlabs.cloud/mcp - Provides 11 database tools via JSON-RPC protocol
• Ollama LLM: all-mpnet-base-v2 model - Generates 768-dimensional embeddings for query similarity
• MCP Protocol (red arrows): Main feature - Dynamic tool discovery, not hardcoded
• SKILL.md Files (orange arrows): 25+ expert guides from cockroachlabs/cockroach GitHub repo
• Query Learning (purple arrows): Optimization - Pre-loads skills from past similar queries (30-40% token savings)"""

    legend_box.text_frame.text = legend_text
    for p in legend_box.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Save
    filename = '/var/www/ai/aidemo2/docs/timbobarchitectureversion2.pptx'
    prs.save(filename)
    print(f"Architecture slide created: {filename}")
    print(f"Total slides: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    create_architecture_slide()
