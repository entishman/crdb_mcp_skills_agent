#!/usr/bin/env python3
"""
Final Visual PowerPoint Presentation - Version 3
Focus: MCP protocol + SKILL.md files for CockroachDB Cloud management
Learning is just an optimization, not the main feature
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def add_box(slide, x, y, width, height, text, fill_color, text_size=10, bold=False):
    """Add a colored box with BLACK text"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(2)

    box.text_frame.text = text
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = 1

    # Set ALL paragraphs to black (multi-line text creates multiple paragraphs)
    for paragraph in box.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(text_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # BLACK TEXT!

    return box

def add_arrow(slide, x1, y1, x2, y2, label=None):
    """Add an arrow connector"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(2)
    connector.line.end_arrow_type = 2

    if label:
        label_box = slide.shapes.add_textbox(
            Inches((x1 + x2) / 2 - 0.4), Inches((y1 + y2) / 2 - 0.15),
            Inches(0.8), Inches(0.25)
        )
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.italic = True
        label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        label_box.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return connector

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = section_title

    p = title_frame.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_bullet_slide(prs, title, bullets):
    """Add a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    tf = body_shape.text_frame
    tf.clear()

    for bullet in bullets:
        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 - level * 2)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(10)

    return slide

def add_architecture_overview(prs):
    """MCP-centric architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title.text = "Architecture: MCP Protocol + SKILL.md Knowledge Base"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # User
    add_box(slide, 4.25, 0.8, 1.5, 0.5, "User", RGBColor(173, 216, 230))
    add_arrow(slide, 5, 1.3, 5, 1.8)

    # Agent (center)
    add_box(slide, 3.5, 1.8, 3, 0.7, "Agent (agent.py)\nOrchestrator", RGBColor(255, 200, 100), 13)

    # Left - Skill System
    add_box(slide, 0.5, 3, 2.5, 0.7, "SKILL.md System\n25+ Best Practice Guides\nStored in ai_demo.skills",
            RGBColor(255, 240, 200), 10)
    add_arrow(slide, 1.75, 2.5, 4.5, 2.2, "fetch guides")

    # Right - Claude API
    add_box(slide, 7, 3, 2.5, 0.7, "Claude Sonnet 4.6\n(Vertex AI)\nTool Use API",
            RGBColor(200, 230, 255), 11)
    add_arrow(slide, 6.5, 2.2, 8, 3, "API calls")

    # Bottom left - Query Learning (OPTIMIZATION)
    add_box(slide, 0.5, 4.2, 2.5, 0.6, "Query Learning\n(OPTIMIZATION)\n768-dim embeddings",
            RGBColor(220, 255, 220), 10)
    add_arrow(slide, 1.75, 3.7, 1.75, 4.2)

    # Bottom center - MCP Client (MAIN FEATURE)
    add_box(slide, 3.5, 4.2, 2.5, 0.6, "MCP Client\n(MAIN PROTOCOL)\nDynamic tool discovery",
            RGBColor(255, 210, 210), 10, True)
    add_arrow(slide, 5, 2.5, 5, 4.2, "execute")

    # Bottom right - MCP Server
    add_box(slide, 6.5, 4.2, 2.5, 0.6, "MCP Server\ncockroachlabs.cloud/mcp\n11 Database Tools",
            RGBColor(240, 220, 255), 10)
    add_arrow(slide, 6.5, 2.5, 7.75, 4.2, "discover")

    # CockroachDB Cluster (THE TARGET)
    add_box(slide, 2.5, 5.5, 5, 0.9, "CockroachDB Cluster\n(THE CLUSTER BEING MANAGED)\nAlso stores: skills, query history, embeddings",
            RGBColor(100, 150, 200), 11, True)
    add_arrow(slide, 5, 4.8, 5, 5.5)

    # Key message
    msg = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.6))
    msg.text_frame.text = "🎯 Core Value: MCP protocol provides 11 database tools + SKILL.md files provide expert guidance\n💡 Optimization: Query learning pre-loads relevant skills (saves 30-40% tokens)"
    for p in msg.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_flow_mcp_basic(prs):
    """Basic MCP flow - tool discovery and execution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 1: MCP Tool Discovery & Execution (Core Feature)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    add_box(slide, 0.3, 0.7, 2.2, 0.55, 'Step 1: User Request\n"list all databases"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    # Step 2
    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: MCP Discovery\n' +
            'mcp_client.py: list_tools()\n' +
            'JSON-RPC: tools/list\n' +
            'MCP Server returns 11 tools',
            RGBColor(255, 210, 210), 9, True)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3
    add_box(slide, 5.6, 0.7, 2.9, 0.55,
            'Step 3: Tool Definitions\n' +
            'list_databases, create_database,\n' +
            'select_query, get_cluster, etc.\n' +
            'Converted to Claude API format',
            RGBColor(240, 220, 255), 9, False)
    add_arrow(slide, 7, 1.25, 7, 1.5)

    # Step 4
    add_box(slide, 5.6, 1.5, 2.9, 0.55,
            'Step 4: Call Claude API\n' +
            'agent.py: run_task()\n' +
            'claude.messages.create(\n' +
            '  tools=[11 MCP tools])',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 7, 2.05, 7, 2.3)

    # Step 5
    add_box(slide, 5.6, 2.3, 2.9, 0.7,
            'Step 5: Claude Decides\n' +
            'Claude reads tool descriptions\n' +
            'Decides: need list_databases\n' +
            'Returns: tool_use block with\n' +
            'tool_name="list_databases"',
            RGBColor(220, 255, 220), 8.5, False)
    add_arrow(slide, 7, 3, 7, 3.25)

    # Step 6
    add_box(slide, 5.6, 3.25, 2.9, 0.75,
            'Step 6: Execute via MCP\n' +
            'mcp_client.py:\n' +
            'execute_tool_call(\n' +
            '  "list_databases", {})\n' +
            'JSON-RPC → MCP Server\n' +
            'MCP → CockroachDB cluster',
            RGBColor(255, 210, 210), 8.5, True)
    add_arrow(slide, 7, 4, 7, 4.25)

    # Step 7
    add_box(slide, 5.6, 4.25, 2.9, 0.65,
            'Step 7: Database Results\n' +
            'CockroachDB: SHOW DATABASES\n' +
            'Returns: defaultdb, ai_demo,\n' +
            'postgres, system\n' +
            'MCP → Agent → Claude',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 5.6, 4.55, 5.3, 4.55)

    # Step 8
    add_box(slide, 2.8, 4.25, 2.5, 0.65,
            'Step 8: Claude Formats\n' +
            'Claude receives data\n' +
            'Formats nice answer\n' +
            'Returns to user',
            RGBColor(255, 240, 200), 9, False)
    add_arrow(slide, 2.8, 4.55, 2.5, 4.55)

    # Step 9
    add_box(slide, 0.3, 4.25, 2.2, 0.65,
            'Step 9: User Sees Answer\n' +
            '"Your cluster has 4\n' +
            'databases: defaultdb,\n' +
            'ai_demo, postgres, system"',
            RGBColor(173, 216, 230), 9, False)

    # Explanations
    explain1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.15), Inches(4.5), Inches(0.55))
    explain1.text_frame.text = "🔑 MCP Protocol: Model Context Protocol provides standardized interface to CockroachDB Cloud. Agent discovers tools dynamically - no hardcoded assumptions!"
    for p in explain1.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)
        p.font.bold = True

    explain2 = slide.shapes.add_textbox(Inches(5), Inches(5.15), Inches(4.5), Inches(0.55))
    explain2.text_frame.text = "🛠️ 11 MCP Tools: list_databases, create_database, list_tables, create_table, select_query, insert_rows, get_cluster, get_cluster_cert, and 3 more. All discovered automatically!"
    for p in explain2.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(5.85), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "💡 This flow works WITHOUT any SKILL.md files. Claude knows basic SQL from training. Tools provide the interface to execute commands."
    for p in legend.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)

    return slide

def add_flow_skill_md(prs):
    """SKILL.md system - the expertise layer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 2: SKILL.md Files - Expert Guidance for Complex Tasks"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    add_box(slide, 0.3, 0.7, 2.2, 0.6, 'Step 1: Complex Request\n"Set up MOLT Replicator\nfor MySQL migration"',
            RGBColor(173, 216, 230), 10, True)
    add_arrow(slide, 2.5, 1, 2.8, 1)

    # Step 2
    add_box(slide, 2.8, 0.7, 2.5, 0.6,
            'Step 2: Claude Realizes\n' +
            '"I need expert guidance\n' +
            'for MOLT migrations"\n' +
            'Has tool: fetch_skill',
            RGBColor(200, 230, 255), 9, False)
    add_arrow(slide, 5.3, 1, 5.6, 1)

    # Step 3
    add_box(slide, 5.6, 0.7, 2.9, 0.6,
            'Step 3: Fetch SKILL.md\n' +
            'Claude calls:\n' +
            'fetch_skill("onboarding-and-\n' +
            'migrations/molt-replicator")',
            RGBColor(255, 240, 200), 8.5, True)
    add_arrow(slide, 7, 1.3, 7, 1.55)

    # Step 4
    add_box(slide, 5.6, 1.55, 2.9, 0.7,
            'Step 4: Load from Database\n' +
            'skill_fetcher.py: get_skill()\n' +
            'Query: ai_demo.skills table\n' +
            'Returns: 8KB best practices doc\n' +
            'From CockroachDB Docs GitHub',
            RGBColor(180, 220, 180), 8.5, False)
    add_arrow(slide, 7, 2.25, 7, 2.5)

    # Step 5
    add_box(slide, 5.6, 2.5, 2.9, 0.65,
            'Step 5: SKILL.md Content\n' +
            '• Prerequisites\n' +
            '• Setup steps\n' +
            '• Common issues\n' +
            '• Best practices',
            RGBColor(255, 240, 200), 8.5, False)
    add_arrow(slide, 7, 3.15, 7, 3.4)

    # Step 6
    add_box(slide, 5.6, 3.4, 2.9, 0.6,
            'Step 6: Expert Analysis\n' +
            'Claude uses SKILL.md to:\n' +
            '1. Check existing setup\n' +
            '2. Guide configuration\n' +
            '3. Provide best practices',
            RGBColor(220, 255, 220), 8.5, False)
    add_arrow(slide, 5.6, 3.7, 5.3, 3.7)

    # Step 7
    add_box(slide, 2.8, 3.4, 2.5, 0.6,
            'Step 7: Use MCP Tools\n' +
            'select_query() to check\n' +
            'create_database() if needed\n' +
            'Based on SKILL.md guidance',
            RGBColor(255, 210, 210), 8.5, False)
    add_arrow(slide, 2.8, 3.7, 2.5, 3.7)

    # Step 8
    add_box(slide, 0.3, 3.4, 2.2, 0.6,
            'Step 8: Complete Guide\n' +
            'Claude provides:\n' +
            '• Step-by-step plan\n' +
            '• Commands to run\n' +
            '• Troubleshooting tips',
            RGBColor(255, 240, 200), 8.5, False)

    # Step 9
    add_box(slide, 0.3, 4.3, 2.2, 0.5,
            'Step 9: Store Learning\n' +
            'query + skills_used\n' +
            'Next time: pre-loaded!',
            RGBColor(220, 255, 220), 9, False)
    add_arrow(slide, 2.5, 4.55, 2.8, 4.55)

    # Step 10
    add_box(slide, 2.8, 4.3, 2.5, 0.5,
            'Step 10: Optimization\n' +
            'Future similar queries\n' +
            'get skill pre-loaded\n' +
            '(saves 1 API call)',
            RGBColor(220, 255, 220), 8.5, False)

    # Explanations
    explain1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.15), Inches(4.5), Inches(0.55))
    explain1.text_frame.text = "📚 SKILL.md Files: 25+ expert guides from cockroachlabs/cockroach GitHub. Cover: migrations, performance, security, operations, troubleshooting. Stored in CockroachDB itself!"
    for p in explain1.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)
        p.font.bold = True

    explain2 = slide.shapes.add_textbox(Inches(5), Inches(5.15), Inches(4.5), Inches(0.55))
    explain2.text_frame.text = "🎯 fetch_skill Tool: Available to Claude just like list_databases. Claude decides when it needs expertise. Not always needed - simple queries work without any skills!"
    for p in explain2.text_frame.paragraphs:
        p.font.size = Pt(8.5)
        p.font.color.rgb = RGBColor(0, 100, 0)

    # Legend
    legend = slide.shapes.add_textbox(Inches(0.3), Inches(5.85), Inches(9.4), Inches(0.3))
    legend.text_frame.text = "🚀 MCP Tools + SKILL.md = Complete Solution. Tools execute operations. Skills provide expert guidance. Together = powerful cluster management assistant."
    for p in legend.text_frame.paragraphs:
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)

    return slide

def add_flow_learning_optimization(prs):
    """Learning as optimization"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.4))
    title.text = "Flow 3: Query Learning - Optimization (Not the Main Feature)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Step 1
    add_box(slide, 0.3, 0.7, 2.2, 0.55, 'Step 1: User Query\n"Check for hot ranges"',
            RGBColor(173, 216, 230), 11, True)
    add_arrow(slide, 2.5, 0.95, 2.8, 0.95)

    # Step 2
    add_box(slide, 2.8, 0.7, 2.5, 0.55,
            'Step 2: Generate Embedding\n' +
            'OPTIONAL OPTIMIZATION\n' +
            'ollama_embedding_manager\n' +
            '→ 768-dim vector',
            RGBColor(255, 230, 230), 9, False)
    add_arrow(slide, 5.3, 0.95, 5.6, 0.95)

    # Step 3
    add_box(slide, 5.6, 0.7, 2.9, 0.6,
            'Step 3: Search History\n' +
            'query_store.py:\n' +
            'find_similar_queries()\n' +
            'CockroachDB pgvector search\n' +
            '✅ Found: "analyze hotspots"',
            RGBColor(180, 255, 180), 8.5, False)
    add_arrow(slide, 7, 1.3, 7, 1.55)

    # Step 4
    add_box(slide, 5.6, 1.55, 2.9, 0.6,
            'Step 4: Pre-load Skills\n' +
            'Past query used:\n' +
            '"observability-and-diagnostics/\n' +
            'analyzing-range-distribution"\n' +
            'Load it before calling Claude',
            RGBColor(255, 240, 200), 8.5, False)
    add_arrow(slide, 7, 2.15, 7, 2.4)

    # Benefits box
    add_box(slide, 5.6, 2.4, 2.9, 0.9,
            'Benefits:\n' +
            '✅ Saves 1 API call\n' +
            '✅ 30-40% fewer tokens\n' +
            '✅ Faster response\n' +
            '✅ Same quality answer',
            RGBColor(220, 255, 220), 9, True)

    # Without optimization
    add_box(slide, 0.3, 1.8, 4, 0.6,
            'WITHOUT Optimization:\n' +
            'User query → Claude → fetch_skill → Claude again\n' +
            '(2 Claude API calls)',
            RGBColor(255, 200, 200), 9, False)
    add_arrow(slide, 4.3, 2.1, 4.6, 2.1)

    # With optimization
    add_box(slide, 4.6, 1.8, 4, 0.6,
            'WITH Optimization:\n' +
            'User query → Pre-load skill → Claude once\n' +
            '(1 Claude API call - faster & cheaper)',
            RGBColor(180, 255, 180), 9, True)

    # How it works
    add_box(slide, 0.3, 3.7, 9.2, 0.7,
            'How Query Learning Works:\n' +
            '1. Every query is converted to 768-dim embedding (sentence-transformers)\n' +
            '2. Stored in CockroachDB with which skills helped answer it\n' +
            '3. Next similar query: pgvector finds match in <1ms (HNSW index)\n' +
            '4. Pre-load those skills → Skip fetch_skill step → Faster & cheaper',
            RGBColor(240, 240, 240), 9, False)

    # Key point
    key = slide.shapes.add_textbox(Inches(0.3), Inches(4.7), Inches(9.2), Inches(0.7))
    key.text_frame.text = "🎯 KEY POINT: Learning is just an OPTIMIZATION\n• Demo works perfectly WITHOUT learning (just slower)\n• Main value: MCP tools + SKILL.md files\n• Learning: Means to efficiently use SKILL.md files"
    for p in key.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 0, 102)

    # Comparison table
    table_title = slide.shapes.add_textbox(Inches(0.3), Inches(5.6), Inches(9.2), Inches(0.25))
    table_title.text_frame.text = "Performance Comparison:"
    for p in table_title.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.3, 5.95, 3, 0.5, "Without Learning\n~3 sec, 2 API calls", RGBColor(255, 200, 200), 10, False)
    add_box(slide, 3.5, 5.95, 3, 0.5, "With Learning (1st time)\n~3 sec, 2 API calls, +store", RGBColor(255, 230, 180), 10, False)
    add_box(slide, 6.7, 5.95, 2.8, 0.5, "With Learning (2nd time)\n~2 sec, 1 API call ✅", RGBColor(180, 255, 180), 10, True)

    return slide

def add_real_world_use_cases(prs):
    """MCP-focused use cases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Real-World Use Cases: Managing CockroachDB Cloud"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Use cases
    add_box(slide, 0.5, 1.2, 4.2, 0.6, "🔧 Cluster Operations", RGBColor(200, 230, 255), 16, True)
    uc1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.2), Inches(1.2))
    uc1.text_frame.text = """• "Show me cluster info"
• "List all databases"
• "Check cluster health"
• Uses: get_cluster, list_databases
• SKILL.md: operations guides"""
    for p in uc1.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.3, 1.2, 4.2, 0.6, "📊 Database Creation", RGBColor(220, 255, 220), 16, True)
    uc2 = slide.shapes.add_textbox(Inches(5.3), Inches(1.9), Inches(4.2), Inches(1.2))
    uc2.text_frame.text = """• "Create test_db for staging"
• "Create users table with schema"
• Uses: create_database, create_table
• With safety confirmations
• SKILL.md: schema design best practices"""
    for p in uc2.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.5, 3.4, 4.2, 0.6, "🔍 Query & Analysis", RGBColor(255, 230, 230), 16, True)
    uc3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(4.2), Inches(1.2))
    uc3.text_frame.text = """• "Show slow queries"
• "Analyze range distribution"
• "Find table hotspots"
• Uses: select_query with expert SQL
• SKILL.md: observability guides"""
    for p in uc3.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.3, 3.4, 4.2, 0.6, "🚀 Migrations", RGBColor(255, 240, 200), 16, True)
    uc4 = slide.shapes.add_textbox(Inches(5.3), Inches(4.1), Inches(4.2), Inches(1.2))
    uc4.text_frame.text = """• "Set up MOLT Replicator"
• "Migrate from MySQL"
• "Best practices for MOLT Fetch"
• Uses: MCP tools + extensive SKILL.md
• Guides through complex processes"""
    for p in uc4.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 2.9, 5.6, 4.2, 0.6, "🎓 Learning & Support", RGBColor(240, 220, 255), 16, True)
    uc5 = slide.shapes.add_textbox(Inches(2.9), Inches(6.3), Inches(4.2), Inches(0.9))
    uc5.text_frame.text = """• "How do I optimize this query?"
• "Security best practices?"
• Pulls from 25+ SKILL.md files
• Consistent expert answers 24/7"""
    for p in uc5.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    return slide

def add_why_cockroachdb(prs):
    """Why CockroachDB - it's the cluster being managed!"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title.text = "Why CockroachDB: The Cluster Being Managed"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Main message
    main = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.7))
    main.text_frame.text = "This demo manages a CockroachDB Cloud cluster using MCP protocol.\nThe same database ALSO stores the agent's metadata (skills, query history)."
    for p in main.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 0, 102)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Features
    add_box(slide, 0.5, 2, 4.3, 0.7, "🎯 Target Cluster", RGBColor(255, 210, 210), 14, True)
    f1 = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.3), Inches(1.3))
    f1.text_frame.text = """• The cluster being managed
• MCP provides tools to interact
• Create databases, tables
• Run queries and analysis
• Monitor health and performance"""
    for p in f1.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.2, 2, 4.3, 0.7, "📚 Metadata Storage", RGBColor(255, 240, 200), 14, True)
    f2 = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4.3), Inches(1.3))
    f2.text_frame.text = """• ai_demo.skills: 25+ SKILL.md files
• ai_demo.query_history: Learning data
• Embeddings with pgvector
• All in the same cluster!
• Eating its own dog food"""
    for p in f2.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 0.5, 4.4, 4.3, 0.7, "🔒 Production Features", RGBColor(220, 255, 220), 14, True)
    f3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(4.3), Inches(1.3))
    f3.text_frame.text = """• ACID transactions for writes
• Strong consistency guarantees
• OAuth authentication
• Distributed globally
• Horizontal scalability"""
    for p in f3.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    add_box(slide, 5.2, 4.4, 4.3, 0.7, "🧠 Vector Capabilities", RGBColor(200, 230, 255), 14, True)
    f4 = slide.shapes.add_textbox(Inches(5.2), Inches(5.2), Inches(4.3), Inches(1.3))
    f4.text_frame.text = """• Native pgvector support
• 768-dim embeddings
• HNSW index for fast search
• No external vector DB needed
• SQL + vectors in one place"""
    for p in f4.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Bottom message
    bottom = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.6))
    bottom.text_frame.text = "🎉 Unique Value: CockroachDB is both the TARGET (cluster being managed) and the INFRASTRUCTURE (stores skills, learning data). One distributed database for everything!"
    for p in bottom.text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return slide

def create_presentation():
    """Create complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(prs, "CockroachDB Cloud Management Agent",
                   "MCP Protocol + SKILL.md Knowledge Base")

    # 2. Overview - MCP FIRST!
    overview = [
        "🎯 What This Demo Does",
        ("Manages CockroachDB Cloud cluster via MCP protocol", 1),
        ("Create databases/tables, run queries, monitor health, all via natural language", 2),
        ("11 database tools discovered dynamically", 1),
        ("Agent learns available operations from MCP server at startup - not hardcoded", 2),
        ("25+ SKILL.md files provide expert guidance", 1),
        ("Best practices for migrations, performance tuning, security, troubleshooting", 2),
        ("Query learning optimizes skill selection", 1),
        ("Pre-loads relevant skills from past similar queries - saves time and tokens", 2),
        "",
        "🔑 Core Components",
        ("MCP Protocol: Connects to CockroachDB Cloud", 1),
        ("Standard protocol eliminates custom integration - secure, reliable communication", 2),
        ("Dynamic Tool Discovery: list_databases, create_table, etc.", 1),
        ("Agent adapts automatically when new database features are added - future-proof", 2),
        ("SKILL.md System: Best practices from GitHub docs", 1),
        ("Official CockroachDB expertise built in - always up-to-date with latest docs", 2),
        ("fetch_skill: Claude calls when expertise needed", 1),
        ("Claude intelligently requests guidance only for complex tasks - efficient context use", 2),
        ("Query Learning: OPTIMIZATION - pre-loads relevant skills", 1),
        ("Remembers which skills helped before - skips fetch_skill step for faster responses", 2),
        "",
        "💡 Value Proposition",
        ("Natural language interface to complex database operations", 1),
        ("Expert guidance from official documentation", 1),
        ("Gets more efficient over time (optional optimization)", 1),
    ]
    add_bullet_slide(prs, "Overview", overview)

    # 3. How to Use the Demo
    setup = [
        "📋 Prerequisites",
        ("CockroachDB Cloud cluster (free tier works!)", 1),
        ("Cluster ID and API key from CockroachDB Cloud", 1),
        ("Google Cloud project with Vertex AI API enabled", 1),
        ("Python 3.8+ with pip", 1),
        "",
        "⚙️ Setup Steps",
        ("1. Copy config.template.json to config.json", 1),
        ("2. Add your CockroachDB cluster_id and api_key", 1),
        ("3. Add your Google Cloud project_id and region", 1),
        ("4. Install dependencies: pip install -r requirements.txt", 1),
        ("5. Run setup: python setup_database.py", 1),
        ("   - Creates ai_demo database", 1),
        ("   - Creates query_history and skills tables", 1),
        ("   - Loads 25+ SKILL.md files from GitHub", 1),
        "",
        "🚀 Running the Demo",
        ("python main.py (interactive mode)", 1),
        ('Ask questions: "list databases", "check cluster health"', 1),
        ("Type 'quit' to exit", 1),
    ]
    add_bullet_slide(prs, "How to Use This Demo", setup)

    # 4. Section: Architecture
    add_section_slide(prs, "Architecture")

    # 5. Architecture diagram
    add_architecture_overview(prs)

    # 6-8. DETAILED FLOW DIAGRAMS - MCP FOCUSED!
    add_flow_mcp_basic(prs)
    add_flow_skill_md(prs)
    add_flow_learning_optimization(prs)

    # 9. Section: Use Cases
    add_section_slide(prs, "Real-World Applications")

    # 10. Use cases
    add_real_world_use_cases(prs)

    # 11. Why CockroachDB
    add_why_cockroachdb(prs)

    # 12. Section: Technical
    add_section_slide(prs, "Technical Details")

    # 13. MCP Protocol details
    mcp_details = [
        "🔌 MCP Protocol (Model Context Protocol)",
        ("Standard protocol for AI-to-database communication", 1),
        ("JSON-RPC 2.0 over HTTPS", 1),
        ("Server: cockroachlabs.cloud/mcp", 1),
        ("Tools/list: Discovers available operations", 1),
        ("Tools/call: Executes database operations", 1),
        "",
        "🛠️ 11 Discovered MCP Tools",
        ("Read: list_databases, list_tables, select_query", 1),
        ("Write: create_database, create_table, insert_rows", 1),
        ("Cluster: get_cluster, get_cluster_cert, get_cluster_api_key", 1),
        ("Plus 2 more... all discovered at runtime!", 1),
        "",
        "📚 SKILL.md Knowledge Base",
        ("25+ expert guides from cockroachlabs/cockroach GitHub", 1),
        ("Categories: operations, security, migrations, diagnostics", 1),
        ("Stored in CockroachDB ai_demo.skills table", 1),
        ("fetch_skill tool: Claude loads when needed", 1),
    ]
    add_bullet_slide(prs, "MCP Protocol & SKILL.md System", mcp_details)

    # 14. Learning optimization
    learning_details = [
        "⚡ Query Learning: Optional Optimization",
        ("NOT the main feature - just makes it faster!", 1),
        ("Works perfectly without learning (just 1 extra API call)", 1),
        "",
        "How It Works:",
        ("Every query → 768-dim embedding (sentence-transformers)", 1),
        ("Store: query text + embedding + skills_used[]", 1),
        ("Next time: pgvector finds similar (HNSW index)", 1),
        ("Pre-load those skills → Skip fetch_skill → Faster", 1),
        "",
        "Performance Impact:",
        ("Without learning: ~3 sec (2 Claude API calls)", 1),
        ("With learning (1st): ~3 sec + store", 1),
        ("With learning (2nd+): ~2 sec (1 Claude API call)", 1),
        ("Saves: 30-40% tokens, 1 API call, ~1 second", 1),
        "",
        "Why It Matters:",
        ("Means to efficiently use SKILL.md files", 1),
        ("Not the goal - just an optimization", 1),
        ("Demo value: MCP + SKILL.md, not learning", 1),
    ]
    add_bullet_slide(prs, "Query Learning (Optimization)", learning_details)

    # 15. Write operations
    write_ops = [
        "✍️ Write Operations (April 5, 2026)",
        ("create_database: Create new databases", 1),
        ("create_table: Create tables with DDL", 1),
        ("insert_rows: Insert data safely", 1),
        "",
        "🔒 Safety Features",
        ("User confirmation REQUIRED for every write", 1),
        ("Shows exact operation before execution", 1),
        ("Can cancel any operation", 1),
        ("OAuth authentication", 1),
        ("Dark orange warning for visibility", 1),
        "",
        "❌ Not Supported (By Design)",
        ("UPDATE, DELETE - too risky", 1),
        ("DROP, TRUNCATE - destructive", 1),
        ("All dangerous operations blocked", 1),
    ]
    add_bullet_slide(prs, "Write Operations with Safety", write_ops)

    # 16. Summary
    summary = [
        "🎯 Core Demo Value",
        ("MCP Protocol: Standard interface to CockroachDB Cloud", 1),
        ("SKILL.md Files: 25+ expert guides from official docs", 1),
        ("Natural language → Database operations", 1),
        ("Complex tasks guided by best practices", 1),
        "",
        "⚡ Optimizations",
        ("Dynamic tool discovery: Not brittle, adapts automatically", 1),
        ("Query learning: Pre-loads relevant skills (30-40% faster)", 1),
        ("Zero-skill mode: Simple queries need no SKILL.md", 1),
        "",
        "🚀 Why CockroachDB",
        ("It's the cluster being managed (via MCP)", 1),
        ("Also stores: SKILL.md files, query history, embeddings", 1),
        ("One database for everything: target + infrastructure", 1),
        ("pgvector: No external vector DB needed", 1),
        ("ACID transactions: Safe writes, no lost learning", 1),
        ("Global distribution: Deploy agents worldwide", 1),
    ]
    add_bullet_slide(prs, "Summary: MCP + SKILL.md = Powerful Agent", summary)

    # Save
    filename = '/var/www/ai/aidemo2/docs/timbobfinalversion3.pptx'
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    create_presentation()
